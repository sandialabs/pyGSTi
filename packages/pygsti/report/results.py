from __future__ import division, print_function, absolute_import, unicode_literals
#*****************************************************************
#    pyGSTi 0.9:  Copyright 2015 Sandia Corporation
#    This Software is released under the GPL license detailed
#    in the file "license.txt" in the top-level pyGSTi directory
#*****************************************************************
""" Defines the Results class and supporting functionality."""

import os  as _os
import re  as _re
import time as _time
import subprocess  as _subprocess
import collections as _collections
import matplotlib  as _matplotlib
import itertools   as _itertools
import copy as _copy

from ..             import objects              as _objs
from ..objects      import gatestring           as _gs
from ..objects      import VerbosityPrinter
from ..construction import spamspecconstruction as _ssc
from ..algorithms   import gaugeopt_to_target   as _optimizeGauge
from ..algorithms   import contract             as _contract
from ..tools        import listtools            as _lt
from ..             import _version

from . import latex      as _latex
from . import generation as _generation
from . import plotting   as _plotting

from .resultcache import ResultCache as _ResultCache

class Results(object):
    """
    Encapsulates a set of GST results.

    A Results object is constructed from the input parameters and raw output
    gatesets from a GST calculation, typically performed by one of the
    "do<something>" methods of GST.Core, and acts as a end-output factory
    (creating reports, presentations, etc), and a derived-results cache
    (so derived quantities don't need to be recomputed many times for
    different output formats).
    """

    def __init__(self, templatePath=None, latexCmd="pdflatex"):
        """
        Initialize a Results object.

        Parameters
        ----------
        templatePath : string or None, optional
            A local path to the stored GST report template files.  The
            default value of None means to use the default path, which
            is almost always what you want.

        latexCmd : string or None, optional
            The system command used to compile latex documents.
        """

        # Internal Flags
        self._bEssentialResultsSet = False
        self._LsAndGermInfoSet = False

        # MPI communicator to be used for confidence region construction
        self._comm = None #TODO: allow this to be passed to __init__?

        # Confidence regions: key == confidence level, val = ConfidenceRegion
        self._confidence_regions = {} # plain dict. Key == confidence level
        self._specials = _ResultCache(self._get_special_fns(), self, "special")

        self.tables  = _ResultCache(self._get_table_fns(), self, "table")
        self.figures = _ResultCache(self._get_figure_fns(), self, "figure")
        #self.qtys = _ResultCache(self._get_qty_fns(), self, "computable qty")

        #Public API parameters
        self.gatesets = {}
        self.gatestring_lists = {}
        self.dataset = None
        self.parameters = {}
        self.options = ResultOptions()
        self.confidence_level = None #holds "current" (i.e. "last")

        # Set default display options (affect how results are displayed)
        self.options.long_tables = False
        self.options.table_class = "pygstiTbl"
        self.options.template_path = templatePath
        self.options.latex_cmd = latexCmd

        # Set default parameter values
        self.parameters = { 'objective': None,
                            'weights':None,
                            'minProbClip': 1e-6,
                            'minProbClipForWeighting': 1e-4,
                            'probClipInterval': (-1e6,1e6),
                            'radius': 1e-4,
                            'hessianProjection': 'optimal gate CIs',
                            'defaultDirectory': None,
                            'defaultBasename': None,
                            'linlogPercentile':  5,
                            'memLimit': None,
                            'gaugeOptParams': {},
                            'cptpPenaltyFactor': 0,
                            'distributeMethod': "deriv" }


    def init_single(self, objective, targetGateset, dataset, gatesetEstimate,
                    gatestring_list, gatesetEstimate_noGaugeOpt=None):
        """
        Initialize this Results object from the inputs and outputs of a
        single (non-iterative) GST method.

        Parameters
        ----------
        objective : {'chi2', 'logl'}
            Whether gateset was obtained by minimizing chi^2 or
            maximizing the log-likelihood.

        targetGateset : GateSet
            The target gateset used when optimizing the objective.

        dataset : DataSet
            The dataset used when optimizing the objective.

        gatesetEstimate : GateSet
            The (single) gate set obtained which optimizes the objective.

        gatestring_list : list of GateStrings
            The list of gate strings used to optimize the objective.

        gatesetEstimate_noGaugeOpt : GateSet, optional
            The value of the estimated gate set *before* any gauge
            optimization was performed on it.

        Returns
        -------
        None
        """

        # Set essential info: gateset estimates(s) but no particular
        # structure known about gateStringLists.
        self.gatesets['target'] = targetGateset
        self.gatesets['iteration estimates'] = [ gatesetEstimate ]
        self.gatesets['final estimate'] = gatesetEstimate
        self.gatestring_lists['iteration'] = [ gatestring_list ]
        self.gatestring_lists['final'] = gatestring_list
        self.dataset = dataset
        self.parameters['objective'] = objective

        if gatesetEstimate_noGaugeOpt is not None:
            self.gatesets['iteration estimates pre gauge opt'] = \
                [ gatesetEstimate_noGaugeOpt ]

        self._bEssentialResultsSet = True


    def init_Ls_and_germs(self, objective, targetGateset, dataset,
                          seedGateset, Ls, germs, gatesetsByL, gateStringListByL,
                          prepStrs, effectStrs, truncFn, fidPairs=None,
                          gatesetsByL_noGaugeOpt=None):

        """
        Initialize this Results object from the inputs and outputs of
        an iterative GST method based on gate string lists containing
        germs repeated up to a maximum-L value that increases with
        iteration.

        Parameters
        ----------
        objective : {'chi2', 'logl'}
            Whether gateset was obtained by minimizing chi^2 or
            maximizing the log-likelihood.

        targetGateset : GateSet
            The target gateset used when optimizing the objective.

        dataset : DataSet
            The dataset used when optimizing the objective.

        seedGateset : GateSet
            The initial gateset used to seed the iterative part
            of the objective optimization.  Typically this is
            obtained via LGST.

        Ls : list of ints
            List of maximum-L values used in the iterations.

        germs : list of GateStrings
            List of germ gate strings used in the objective optimization.

        gatesetsByL : list of GateSets
            The estimated gateset at each L value.

        gateStringListByL : list of lists of GateStrings
            The gate string list used at each L value.

        prepStrs : list of GateStrings
            The list of state preparation fiducial strings
            in the objective optimization.

        effectStrs : list of GateStrings
            The list of measurement fiducial strings
            in the objective optimization.

        truncFn : function
            The truncation function used, indicating how a
            germ should be repeated "L times".  Function should
            take parameters (germ, L) and return the repeated
            gate string.  For example, see
            pygsti.construction.repeat_with_max_length.

        fidPairs : list or dict, optional
            Specifies a subset of all prepStr,effectStr string pairs to be used in
            reports.  If `fidPairs` is a list, each element of `fidPairs` is a
            ``(iRhoStr, iEStr)`` 2-tuple of integers, which index a string within
            the state preparation and measurement fiducial strings respectively. If
            `fidPairs` is a dict, then the keys must be germ strings and values are
            lists of 2-tuples as in the previous case.

        gatesetsByL_noGaugeOpt : list of GateSets, optional
            The value of the estimated gate sets *before* any gauge
            optimization was performed on it.

        Returns
        -------
        None
        """

        assert(len(gateStringListByL) == len(gatesetsByL) == len(Ls))

        # Set essential info: gateset estimates(s) but no particular
        # structure known about gateStringLists.
        self.gatesets['target'] = targetGateset
        self.gatesets['seed'] = seedGateset
        self.gatesets['iteration estimates'] = gatesetsByL
        self.gatesets['final estimate'] = gatesetsByL[-1]
        self.gatestring_lists['iteration'] = gateStringListByL
        self.gatestring_lists['final'] = gateStringListByL[-1]
        self.gatestring_lists['all'] = _lt.remove_duplicates(
            list(_itertools.chain(*gateStringListByL)) )

        running_lst = []; delta_lsts = []
        for L,lst in zip(Ls,gateStringListByL):
            delta_lst = [ x for x in lst if (x not in running_lst) ]
            #if L != 0: running_lst += delta_lst # L=0 is special case - doesn't count in running list
            delta_lsts.append(delta_lst)
        self.gatestring_lists['iteration delta'] = delta_lsts # *added* at each iteration

        self.dataset = dataset
        self.parameters['objective'] = objective
        if gatesetsByL_noGaugeOpt is not None:
            self.gatesets['iteration estimates pre gauge opt'] = \
                gatesetsByL_noGaugeOpt

        self._bEssentialResultsSet = True

        #Set "Ls and germs" info: gives particular structure
        # to the gateStringLists used to obtain estimates
        self.gatestring_lists['prep fiducials'] = prepStrs
        self.gatestring_lists['effect fiducials'] = effectStrs
        self.gatestring_lists['germs'] = germs
        self.parameters['max length list'] = Ls
        self.parameters['fiducial pairs'] = fidPairs
        self.parameters['L,germ tuple base string dict'] = \
            _collections.OrderedDict( [ ( (L,germ), truncFn(germ,L) )
                                        for L in Ls for germ in germs] )
        self._LsAndGermInfoSet = True

    def reoptimize_gauge(self, gaugeOptParams, setparam=True):
        """
        Re-optimizes the gauge of the final gateset.

        This function updates the value of this object's 
        `gatesets['final estimate']` gate set with the result of the specified
        gauge optimization, and also clears cached figures, tables, etc. which
        are gauge dependent to that they are re-computed using the updated gate
        set.

        Parameters
        ----------
        gaugeOptParams : dict, optional
            A dictionary of arguments to :func:`gaugeopt_to_target`, specifying
            how the gauge optimization should be performed.  The keys and
            values of this dictionary may correspond to any of the arguments
            of :func:`gaugeopt_to_target` *except* for the first `gateset` 
            argument, which is taken to be `gatesets['final']`.  The 
            `targetGateset` argument *can* be specified, but if it isn't, is
            taken to be `gatesets['target']`.  This argument may also be a
            list of such dictionaries, in which case each element describes
            a successive stage of gauge optimization.

        setparam : bool, optional
            Whether to set `parameters['gaugeOptParams']` to the list of
            parameter dictionaries returned by this function.

        Returns
        -------
        List of OrderedDicts
            A list of dictionaries, each containing gauge optimization
            parameters for a single stage of gauge optimization.
        """
        assert(self._bEssentialResultsSet)

        if hasattr(gaugeOptParams,"keys"):
            go_params_list = [gaugeOptParams]
        else: go_params_list = gaugeOptParams

        ordered_go_params_list = []
        for go_params in go_params_list:
            if "targetGateset" not in go_params:
                go_params["targetGateset"] = self.gatesets['target']

            ordered_go_params_list.append( _collections.OrderedDict( 
                [(k,go_params[k]) for k in sorted(list(go_params.keys()))]))

            self.gatesets['final estimate'] = _optimizeGauge(
                self.gatesets['final estimate'],**go_params)
            
        if setparam:
            self.parameters['gaugeOptParams'] = ordered_go_params_list
            
        #Clear everything that is (possibly) gauge dependent
        #  Note: also clear 'bestGatesetGaugeOptParamsTable' since we might have updated params
        except_tables = ['fiducialListTable', 'prepStrListTable',
                         'effectStrListTable', 'germListTable',
                         'germList2ColTable', 'chi2ProgressTable',
                         'logLProgressTable', 'progressTable',
                         'byGermTable', 'bestGatesetEvalTable']
        except_figures = [ "colorBoxPlotKeyPlot", "bestEstimateColorBoxPlot",
                           "invertedBestEstimateColorBoxPlot",
                           "bestEstimateSummedColorBoxPlot",
                           "blankBoxPlot", "blankSummedBoxPlot"]
        except_specials = [ 'blankGaugeOptAppendixTables',
                            'bestEstimateColorBoxPlotPages']

        if 'max length list' in self.parameters:
            except_figures += ["estimateForLIndex%dColorBoxPlot" % i 
                     for i in range(len(self.parameters['max length list']))]

        self._confidence_regions = {}
        self._specials.clear_cached_data(except_specials)
        self.tables.clear_cached_data(except_tables)
        self.figures.clear_cached_data(except_figures)
        
        return ordered_go_params_list


    def copy(self):
        """ Creates a copy of this Results object. """
        cpy = Results(self.options.template_path, self.options.latex_cmd)
        cpy._bEssentialResultsSet = self._bEssentialResultsSet
        cpy._LsAndGermInfoSet = self._LsAndGermInfoSet
        cpy._comm = self._comm
        cpy._confidence_regions = self._confidence_regions.copy()
        cpy._specials = self._specials.copy()
        cpy.tables = self.tables.copy()
        cpy.figures = self.figures.copy()
        cpy.gatesets = self.gatesets.copy()
        cpy.gatestring_lists = self.gatestring_lists.copy()
        cpy.dataset = self.dataset.copy()
        cpy.parameters = self.parameters.copy()
        cpy.options = self.options.copy()
        cpy.confidence_level = self.confidence_level
        return cpy


    def __getstate__(self):
        #Return the state (for pickling) -- *don't* pickle Comm object
        to_pickle = self.__dict__.copy()
        del to_pickle['_comm'] # one *cannot* pickle Comm objects
        return to_pickle


    def __setstate__(self, stateDict):
        #Must set ResultCache parent & functions, since these are
        # not pickled (to avoid circular pickle references)
        self.__dict__.update(stateDict)
        self._comm = None
        self._specials._setparent(self._get_special_fns(), self)
        self.tables._setparent(self._get_table_fns(), self)
        self.figures._setparent(self._get_figure_fns(), self)

    def __str__(self):
        s  = "----------------------------------------------------------\n"
        s += "---------------- pyGSTi Results Object -------------------\n"
        s += "----------------------------------------------------------\n"
        s += "\n"
        s += "I can create reports for you directly, via my create_XXX\n"
        s += "functions, or you can query me for result data via members:\n\n"
        s += " .dataset    -- the DataSet used to generate these results\n\n"
        s += " .gatesets   -- a dictionary of GateSet objects w/keys:\n"
        s += " ---------------------------------------------------------\n"
        s += "  " + "\n  ".join(list(self.gatesets.keys())) + "\n"
        s += "\n"
        s += " .gatestring_lists   -- a dict of GateString lists w/keys:\n"
        s += " ---------------------------------------------------------\n"
        s += "  " + "\n  ".join(list(self.gatestring_lists.keys())) + "\n"
        s += "\n"
        s += " .tables   -- a dict of ReportTable objects w/keys:\n"
        s += " ---------------------------------------------------------\n"
        s += "  " + "\n  ".join(list(self.tables.keys())) + "\n"
        s += "\n"
        s += " .figures   -- a dict of ReportFigure objects w/keys:\n"
        s += " ---------------------------------------------------------\n"
        s += "  " + "\n  ".join(list(self.figures.keys())) + "\n"
        s += "\n"
        s += " .parameters   -- a dict of simulation parameters:\n"
        s += " ---------------------------------------------------------\n"
        s += "  " + "\n  ".join(list(self.parameters.keys())) + "\n"
        s += "\n"
        s += " .options   -- a container of display options:\n"
        s += " ---------------------------------------------------------\n"
        s += self.options.describe("   ") + "\n"
        s += "\n"
        s += "NOTE: passing 'tips=True' to create_full_report_pdf or\n"
        s += " create_brief_report_pdf will add markup to the resulting\n"
        s += " PDF indicating how tables and figures in the PDF correspond\n"
        s += " to the values of .tables[ ] and .figures[ ] listed above.\n"
        return s

    def _process_call(self, call):
        process = _subprocess.Popen(call, stdout=_subprocess.PIPE,
                                    stderr=_subprocess.PIPE)
        stdout, stderr = process.communicate()
        return stdout, stderr, process.returncode

    def _evaluate_call(self, call, stdout, stderr, returncode, printer):
        if len(stderr) > 0:
            printer.error(stderr)
        if returncode > 0:
            raise _subprocess.CalledProcessError(returncode, call)

    def _compile_latex_report(self, report_dir, report_base, latex_call,
                              printer):
        """Compile a PDF report from a TeX file. Will compile twice
        automatically.

        Parameters
        ----------
        report_dir : string
            The directory for the output file.

        report_base : string
            The base name for the output file (not including any extensions).

        latex_call : list of string
            List containing the command and flags in the form that
            :function:`subprocess.check_call` uses.

        printer : VerbosityPrinter
            Printer to handle logging.

        Raises
        ------
        subprocess.CalledProcessException
            If the call to the process comiling the PDF returns non-zero exit
            status.

        """
        texFilename = report_base + ".tex"
        pdfPathname = _os.path.join(report_dir, report_base + ".pdf")
        call = latex_call + [texFilename]
        stdout, stderr, returncode = self._process_call(call)
        self._evaluate_call(call, stdout, stderr, returncode, printer)
        printer.log("Initial output PDF %s successfully generated." %
                    pdfPathname)
        # We could check if the log file contains "Rerun" in it,
        # but we'll just re-run all the time now
        stdout, stderr, returncode = self._process_call(call)
        self._evaluate_call(call, stdout, stderr, returncode, printer)
        printer.log("Final output PDF %s successfully generated. " %
                    pdfPathname + "Cleaning up .aux and .log files.")
        _os.remove( report_base + ".log" )
        _os.remove( report_base + ".aux" )


    def _get_table_fns(self):
        """
        Return a dictionary of functions which create a table identified by
        the dictionary key.  These functions are used for the lazy creation
        of tables within the "tables" member of a Results instance.
        """

        #Validation functions: return a list of the computable key(s)
        # which match their single "key" argument.  It can be assumed
        # that "key" is either equal to or matches the corresponding
        # compute-function key.  Since the latter may be a regular expression,
        # "key" may also be this same regular-expression, in which case
        # a list of currently computable keys (based on current Results
        # parameters, etc.) should be returned.  In the more mundane
        # case where key is just a string, the function simply returns
        # that same string when that key can be computed, and None
        # otherwise.
        def validate_none(key):
            return [key]
        def validate_essential(key):
            return [key] if self._bEssentialResultsSet else []
        def validate_LsAndGerms(key):
            return [key] if (self._bEssentialResultsSet and
                             self._LsAndGermInfoSet) else []
        def noConfidenceLevelDependence(level):
            """ Designates a table as independent of the confidence level"""
            if level is not None: raise _ResultCache.NoCRDependenceError

        def setup():
            return (self.gatesets['target'], self.gatesets['final estimate'])

        fns = _collections.OrderedDict()

        def fn(key, confidenceLevel, vb):
            return _generation.get_blank_table()
        fns['blankTable'] = (fn, validate_none)

        # target gateset tables
        def fn(key, confidenceLevel, vb):
            gsTgt, _ = setup()
            return _generation.get_gateset_spam_table(gsTgt, None)
        fns['targetSpamTable'] = (fn, validate_essential)


        def fn(key, confidenceLevel, vb):
            gsTgt, _ = setup()
            return _generation.get_gateset_spam_table(gsTgt, None, None, False)
        fns['targetSpamBriefTable'] = (fn, validate_essential)


        def fn(key, confidenceLevel, vb):
            gsTgt, _ = setup()
            return _generation.get_unitary_gateset_gates_table(
                gsTgt, None)
        fns['targetGatesTable'] = (fn, validate_essential)


        # dataset and gatestring list tables
        def fn(key, confidenceLevel, vb):
            #maxLen = max( 2*max( map(len,self.prepStrs + self.effectStrs) ),
            #             10 ) #heuristic (unused)
            gsTgt, _ = setup()
            if self._LsAndGermInfoSet:
                strs = ( self.gatestring_lists['prep fiducials'],
                         self.gatestring_lists['effect fiducials'] )
            else: strs = None
            return _generation.get_dataset_overview_table(
                self.dataset, gsTgt, 10, strs)
        fns['datasetOverviewTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            setup()
            strs = ( self.gatestring_lists['prep fiducials'],
                     self.gatestring_lists['effect fiducials'] )

            return _generation.get_gatestring_multi_table(
                strs, ["Prep.","Measure"], "Fiducials")
        fns['fiducialListTable'] = (fn, validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            return _generation.get_gatestring_table(
                self.gatestring_lists['prep fiducials'],
                "Preparation Fiducial")
        fns['prepStrListTable'] = (fn, validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            return _generation.get_gatestring_table(
                self.gatestring_lists['effect fiducials'],
                "Measurement Fiducial")
        fns['effectStrListTable'] = (fn, validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            return _generation.get_gatestring_table(
                self.gatestring_lists['germs'], "Germ")
        fns['germListTable'] = (fn, validate_LsAndGerms)


        def fn(key, confidenceLevel, vb):
            return _generation.get_gatestring_table(
                self.gatestring_lists['germs'], "Germ", nCols=2)
        fns['germList2ColTable'] = (fn, validate_LsAndGerms)


        # Estimated gateset tables
        def fn(key, confidenceLevel, vb):
            _, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_gateset_spam_table(gsBest, None, cri)
        fns['bestGatesetSpamTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_gateset_spam_table(gsBest, gsTgt, cri, False)
        fns['bestGatesetSpamBriefTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            _, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_gateset_spam_parameters_table(gsBest, cri)
        fns['bestGatesetSpamParametersTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            _, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_gateset_gates_table(gsBest, cri)
        fns['bestGatesetGatesTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            _, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_gateset_choi_table(gsBest, cri)
        fns['bestGatesetChoiTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            _, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_gateset_decomp_table(gsBest, cri)
        fns['bestGatesetDecompTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            _, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_gateset_rotn_axis_table(gsBest, cri, True)
        fns['bestGatesetRotnAxisTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_gateset_eigenval_table(gsBest, gsTgt, cri)
        fns['bestGatesetEvalTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            _, gsBest = setup()
            #cri = self._get_confidence_region(confidenceLevel)
            return _generation.get_gateset_closest_unitary_table(gsBest) #, cri)
        fns['bestGatesetClosestUnitaryTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            #Note: ALWAYS compute error bars if cri is not None
            return _generation.get_gates_vs_target_table(gsBest, gsTgt, cri)
        fns['bestGatesetVsTargetTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, _ = setup()
            noConfidenceLevelDependence(confidenceLevel)
            best_gs_gauges = self._specials.get(
                'singleGateTargetGaugeOptGatesets',verbosity=vb)
            return _generation.get_selected_gates_vs_target_table(
                best_gs_gauges, gsTgt, None)
        fns['gaugeOptGatesetsVsTargetTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            cptp_go_gateset = self._specials.get(
                'CPTPGaugeOptGateset',verbosity=vb)
            return _generation.get_gateset_choi_eigenval_table(
                cptp_go_gateset, "goCPTPChoiEvalBars")
        fns['gaugeOptCPTPGatesetChoiTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_spam_vs_target_table(gsBest, gsTgt, cri)
        fns['bestGatesetSpamVsTargetTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_gates_vs_target_err_gen_table(
                gsBest, gsTgt, cri, self.options.errgen_type)
        fns['bestGatesetErrorGenTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_gates_vs_target_angles_table(
                gsBest, gsTgt, cri)
        fns['bestGatesetVsTargetAnglesTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            setup()
            return _generation.get_gaugeopt_params_table(
                self.parameters['gaugeOptParams'])
        fns['bestGatesetGaugeOptParamsTable'] = (fn, validate_essential)

        # progress tables
        def fn(key, confidenceLevel, vb):
            setup()
            return _generation.get_chi2_progress_table(
                self.parameters['max length list'],
                self.gatesets['iteration estimates'],
                self.gatestring_lists['iteration'], self.dataset)
        fns['chi2ProgressTable'] = (fn, validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            setup()
            return _generation.get_logl_progress_table(
                self.parameters['max length list'],
                self.gatesets['iteration estimates'],
                self.gatestring_lists['iteration'], self.dataset)
        fns['logLProgressTable'] = (fn, validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            setup()
            if self.parameters['objective'] == "logl":
                return _generation.get_logl_progress_table(
                    self.parameters['max length list'],
                    self.gatesets['iteration estimates'],
                    self.gatestring_lists['iteration'], self.dataset)
            elif self.parameters['objective'] == "chi2":
                return _generation.get_chi2_progress_table(
                    self.parameters['max length list'],
                    self.gatesets['iteration estimates'],
                    self.gatestring_lists['iteration'], self.dataset)
            else: raise ValueError("Invalid Objective: %s" %
                                   self.parameters['objective'])
        fns['progressTable'] = (fn, validate_LsAndGerms)


        def fn(key, confidenceLevel, vb):
            #Much of the below setup is similar to plot_setup() in _get_figure_fns
            strs  = (self.gatestring_lists['prep fiducials'],
                     self.gatestring_lists['effect fiducials'])
            germs = [ g for g in self.gatestring_lists['germs'] if len(g) <= 3 ]
            gsBest = self.gatesets['final estimate']
            fidPairs = self.parameters['fiducial pairs']
            Ls = self.parameters['max length list']

            if fidPairs is None: fidpair_filters = None
            elif isinstance(fidPairs,dict) or hasattr(fidPairs,"keys"):
                #Assume fidPairs is a dict indexed by germ
                fidpair_filters = { (x,y): fidPairs[germ] 
                                    for x in Ls for y in germs }
            else:
                #Assume fidPairs is a list
                fidpair_filters = { (x,y): fidPairs
                                    for x in Ls for y in germs }

            gstr_filters = { (x,y) : self.gatestring_lists['iteration'][i]
                             for i,x in enumerate(Ls) for y in germs }

            if self.parameters['objective'] == "logl":
                return _generation.get_logl_bygerm_table(
                    gsBest, self.dataset, germs, strs, Ls,
                    self.parameters['L,germ tuple base string dict'],
                    fidpair_filters, gstr_filters)
            elif self.parameters['objective'] == "chi2":
                raise NotImplementedError("byGermTable not implemented for chi2 objective")
            else: raise ValueError("Invalid Objective: %s" %
                                   self.parameters['objective'])
        fns['byGermTable'] = (fn, validate_LsAndGerms)


        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            cptp_go_gateset = self._specials.get('CPTPGaugeOptGateset',verbosity=vb).copy()
            cptp_go_gateset.set_all_parameterizations("full") #for contraction
            cptp_gateset = _contract(cptp_go_gateset, "CPTP")
            return _generation.get_logl_projected_err_gen_table(
                gsBest, gsTgt, self.gatestring_lists['final'], self.dataset,
                cptp_gateset, self.options.errgen_type)
        fns['logLErrgenProjectionTable'] = (fn, validate_essential)


        # figure-containing tables
        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            return _generation.get_gateset_gate_boxes_table(
                [gsTgt], ["targetGatesBoxes"], maxHeight=4.0)
        fns['targetGatesBoxTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            return _generation.get_gateset_gate_boxes_table(
                [gsTgt, gsBest], ["targetGatesBoxes", "bestGatesBoxes"],
                ['Target','Estimated'])
        fns['bestGatesetGatesBoxTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            return _generation.get_gates_vs_target_err_gen_boxes_table(
                gsBest, gsTgt, "bestErrgenBoxes", genType=self.options.errgen_type)
        fns['bestGatesetErrGenBoxTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            return _generation.get_projected_err_gen_comparison_table(
                gsBest, gsTgt, compare_with='target', genType=self.options.errgen_type)
        fns['bestGatesetErrGenProjectionTargetMetricsTable'] = (fn,validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            return _generation.get_projected_err_gen_comparison_table(
                gsBest, gsTgt, compare_with='estimate', genType=self.options.errgen_type)
        fns['bestGatesetErrGenProjectionSelfMetricsTable'] = (fn,validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            return _generation.get_gateset_eigenval_table(
                gsBest, gsTgt, "bestEvalPolarPlt")
        fns['bestGatesetEvalTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            gsTgt, gsBest = setup()
            return _generation.get_gateset_relative_eigenval_table(
                gsBest, gsTgt, "bestRelEvalPolarPlt",
                genType=self.options.errgen_type)
        fns['bestGatesetRelEvalTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            _, gsBest = setup()
            cri = self._get_confidence_region(confidenceLevel)
            if cri and cri.has_hessian() == False: cri = None
            return _generation.get_gateset_choi_eigenval_table(
                gsBest, "bestChoiEvalBars", confidenceRegionInfo=cri)
        fns['bestGatesetChoiEvalTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            _, gsBest = setup()
            noConfidenceLevelDependence(confidenceLevel)
            return _generation.get_pauli_err_gen_projector_boxes_table(
                gsBest.dim, "hamiltonian", "pauli_ham")
        fns['hamiltonianProjectorTable'] = (fn, validate_essential)

        def fn(key, confidenceLevel, vb):
            _, gsBest = setup()
            noConfidenceLevelDependence(confidenceLevel)
            return _generation.get_pauli_err_gen_projector_boxes_table(
                gsBest.dim, "stochastic", "pauli_sto")
        fns['stochasticProjectorTable'] = (fn, validate_essential)

        return fns


    def _get_figure_fns(self):
        """
        Return a dictionary of functions which create a figure identified by
        the dictionary key.  These functions are used for the lazy creation
        of figures within the "figures" member of a Results instance.
        """

        def getPlotFn():
            obj = self.parameters['objective']
            assert(obj in ("chi2","logl"))
            if   obj == "chi2": return _plotting.chi2_boxplot
            elif obj == "logl": return _plotting.logl_boxplot

        def getDirectPlotFn():
            obj = self.parameters['objective']
            assert(obj in ("chi2","logl"))
            if   obj == "chi2": return _plotting.direct_chi2_boxplot
            elif obj == "logl": return _plotting.direct_logl_boxplot

        def getWhackAMolePlotFn():
            obj = self.parameters['objective']
            assert(obj in ("chi2","logl"))
            if   obj == "chi2": return _plotting.whack_a_chi2_mole_boxplot
            elif obj == "logl": return _plotting.whack_a_logl_mole_boxplot

        def getMPC():
            obj = self.parameters['objective']
            assert(obj in ("chi2","logl"))
            if obj == "chi2":
                return self.parameters['minProbClipForWeighting']
            elif obj == "logl":
                return self.parameters['minProbClip']

        def plot_setup():
            m = 0
            M = 10
            baseStr_dict = self._getBaseStrDict()
            strs  = (self.gatestring_lists['prep fiducials'],
                     self.gatestring_lists['effect fiducials'])
            germs = self.gatestring_lists['germs']
            gsBest = self.gatesets['final estimate']
            fidPairs = self.parameters['fiducial pairs']
            Ls = self.parameters['max length list']
            st = 1 if Ls[0] == 0 else 0 #start index: skip LGST column in plots

            if fidPairs is None: fidpair_filters = None
            elif isinstance(fidPairs,dict) or hasattr(fidPairs,"keys"):
                #Assume fidPairs is a dict indexed by germ
                fidpair_filters = { (x,y): fidPairs[germ] 
                                    for x in Ls[st:] for y in germs }
            else:
                #Assume fidPairs is a list
                fidpair_filters = { (x,y): fidPairs
                                    for x in Ls[st:] for y in germs }

            gstr_filters = { (x,y) : self.gatestring_lists['iteration'][i]
                             for i,x in enumerate(Ls[st:],start=st)
                             for y in germs }
            return Ls,germs,gsBest,fidpair_filters,gstr_filters,m,M,baseStr_dict,strs,st

        def noConfidenceLevelDependence(level):
            """ Designates a figure as independent of the confidence level"""
            if level is not None: raise _ResultCache.NoCRDependenceError

        def validate_LsAndGerms(key):
            return [key] if (self._bEssentialResultsSet and
                             self._LsAndGermInfoSet) else []


        fns = _collections.OrderedDict()


        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            strs  = (self.gatestring_lists['prep fiducials'],
                     self.gatestring_lists['effect fiducials'])
            return _plotting.gof_boxplot_keyplot(strs)
        fns["colorBoxPlotKeyPlot"] = (fn,validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            plotFn = getPlotFn();  mpc = getMPC()
            Ls,germs, gsBest, fpr_filters, gstr_filters, _, _, baseStr_dict, strs, st = plot_setup()
            return plotFn(Ls[st:], germs, baseStr_dict,
                          self.dataset, gsBest, strs,
                          r"$L$", "germ", scale=1.0, sumUp=False,
                          histogram=True, title="", 
                          fidpair_filters=fpr_filters,
                          gatestring_filters = gstr_filters,
                          linlg_pcntle=float(self.parameters['linlogPercentile']) / 100,
                          minProbClipForWeighting=mpc, save_to="", ticSize=20)
        fns["bestEstimateColorBoxPlot"] = (fn,validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            plotFn = getPlotFn(); mpc = getMPC()
            Ls,germs, gsBest, fpr_filters, gstr_filters, _, _, baseStr_dict, strs, st = plot_setup()
            return plotFn( Ls[st:], germs, baseStr_dict,
                           self.dataset, gsBest, strs,
                           r"$L$", "germ", scale=1.0, sumUp=False,
                           histogram=True, title="",
                           fidpair_filters=fpr_filters,
                           gatestring_filters = gstr_filters,
                           linlg_pcntle=float(self.parameters['linlogPercentile']) / 100,
                           save_to="", ticSize=20, minProbClipForWeighting=mpc,
                           invert=True)
        fns["invertedBestEstimateColorBoxPlot"] = (fn,validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            plotFn = getPlotFn();  mpc = getMPC()
            Ls,germs, gsBest, fpr_filters, gstr_filters, _, _, baseStr_dict, strs, st = plot_setup()
            #sumScale = len(strs[0])*len(strs[1]) \
            #    if fidPairs is None else len(fidPairs)
            return plotFn( Ls[st:], germs, baseStr_dict,
                           self.dataset, gsBest, strs,
                           r"$L$", "germ", scale=1.0,
                           sumUp=True, histogram=False, title="",
                           fidpair_filters=fpr_filters,
                           gatestring_filters = gstr_filters,
                           minProbClipForWeighting=mpc,
                           save_to="", ticSize=14, linlg_pcntle=float(self.parameters['linlogPercentile']) / 100)
        fns["bestEstimateSummedColorBoxPlot"] = (fn,validate_LsAndGerms)


        expr1 = "estimateForLIndex(\d+?)ColorBoxPlot"
        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            plotFn = getPlotFn();  mpc = getMPC()
            Ls,germs, _, fpr_filters, gstr_filters, _, _, baseStr_dict, strs, st = plot_setup()
            i = int(_re.match(expr1,key).group(1))
            return plotFn( Ls[st:i+1], germs, baseStr_dict,
                           self.dataset, self.gatesets['iteration estimates'][i],
                           strs, r"$L$", "germ", scale=1.0, sumUp=False,
                           histogram=False, title="",
                           fidpair_filters=fpr_filters,
                           gatestring_filters = gstr_filters,
                           linlg_pcntle=float(self.parameters['linlogPercentile']) / 100,
                           save_to="", minProbClipForWeighting=mpc, ticSize=20)
        def fn_validate(key):
            if not self._LsAndGermInfoSet: return []

            keys = ["estimateForLIndex%dColorBoxPlot" % i
                    for i in range(len(self.parameters['max length list']))]
            if key == expr1: return keys # all computable keys
            elif key in keys: return [key]
            else: return []
        fns[expr1] = (fn, fn_validate)

        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            Ls,germs, _, _, _, _, _, baseStr_dict, strs, st = plot_setup()
            return _plotting.blank_boxplot(
                Ls[st:], germs, baseStr_dict, strs, r"$L$", "germ",
                scale=1.0, title="", sumUp=False, save_to="", ticSize=20)
        fns["blankBoxPlot"] = (fn,validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            Ls, germs, _, _, _, _, _, baseStr_dict, strs, st = plot_setup()
            return _plotting.blank_boxplot(
                Ls[st:], germs, baseStr_dict, strs, r"$L$", "germ",
                scale=1.0, title="", sumUp=True, save_to="", ticSize=20)
        fns["blankSummedBoxPlot"] = (fn,validate_LsAndGerms)


        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            directPlotFn = getDirectPlotFn(); mpc = getMPC()
            Ls, germs, _, _, _, _, _, baseStr_dict, strs, st = plot_setup()
            directLGST = self._specials.get('direct_lgst_gatesets',verbosity=vb)
            return directPlotFn( Ls[st:], germs, baseStr_dict, self.dataset,
                                 directLGST, strs, r"$L$", "germ",
                                 scale=1.0, sumUp=False, title="",
                                 minProbClipForWeighting=mpc,
                                 fidpair_filters= None,
                                 gatestring_filters = None, #don't use filters for direct plots
                                 save_to="", ticSize=20, linlg_pcntle=float(self.parameters['linlogPercentile']) / 100)
        fns["directLGSTColorBoxPlot"] = (fn,validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            directPlotFn = getDirectPlotFn(); mpc = getMPC()
            Ls, germs, _, _, _, _, _, baseStr_dict, strs, st = plot_setup()
            directLongSeqGST = self._specials.get('DirectLongSeqGatesets',
                                                  verbosity=vb)
            return directPlotFn( Ls[st:], germs, baseStr_dict, self.dataset,
                                 directLongSeqGST, strs, r"$L$", "germ",
                                 scale=1.0, sumUp=False, title="",
                                 minProbClipForWeighting=mpc,
                                 fidpair_filters = None,
                                 gatestring_filters = None, #don't use filter_dict for direct plots
                                 save_to="", ticSize=20, linlg_pcntle=float(self.parameters['linlogPercentile']) / 100)
        fns["directLongSeqGSTColorBoxPlot"] = (fn,validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            Ls, germs, gsBest, _, _, _, _, baseStr_dict, _, st = plot_setup()
            directLGST = self._specials.get('direct_lgst_gatesets',verbosity=vb)
            return _plotting.direct_deviation_boxplot(
                Ls[st:], germs, baseStr_dict, self.dataset,
                gsBest, directLGST, r"$L$", "germ", scale=1.0,
                prec=-1, title="", save_to="", ticSize=20)
        fns["directLGSTDeviationColorBoxPlot"] = (fn,validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            Ls, germs, gsBest, _, _, _, _, baseStr_dict, _, st = plot_setup()
            directLongSeqGST = self._specials.get('DirectLongSeqGatesets',
                                                  verbosity=vb)
            return _plotting.direct_deviation_boxplot(
                Ls[st:], germs, baseStr_dict, self.dataset,
                gsBest, directLongSeqGST, r"$L$", "germ",
                scale=1.0, prec=-1, title="", save_to="", ticSize=20)
        fns["directLongSeqGSTDeviationColorBoxPlot"] = (fn,validate_LsAndGerms)

        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            Ls, germs, _, _, _, _, _, baseStr_dict, _, st = plot_setup()
            directLongSeqGST = self._specials.get('DirectLongSeqGatesets',
                                                  verbosity=vb)
            return _plotting.small_eigval_err_rate_boxplot(
                Ls[st:], germs, baseStr_dict, self.dataset,
                directLongSeqGST, r"$L$", "germ", scale=1.0, title="",
                save_to="", ticSize=20)
        fns["smallEigvalErrRateColorBoxPlot"] = (fn,validate_LsAndGerms)


        expr2 = "whack(.+?)MoleBoxes"
        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            Ls,germs, gsBest, fpr_filters, gstr_filters, _, _, baseStr_dict, strs, st = plot_setup()
            highestL = Ls[-1]; hammerWeight = 10.0; mpc = getMPC()
            gateLabel = _re.match(expr2,key).group(1)
            strToWhack = _gs.GateString( (gateLabel,)*highestL )
            whackAMolePlotFn = getWhackAMolePlotFn()
            return whackAMolePlotFn(strToWhack, self.gatestring_lists['all'],
                                    Ls[st:], germs, baseStr_dict, self.dataset,
                                    gsBest, strs, r"$L$", "germ", scale=1.0,
                                    sumUp=False,title="",whackWith=hammerWeight,
                                    save_to="", minProbClipForWeighting=mpc,
                                    ticSize=20, fidpair_filters=fpr_filters,
                                    gatestring_filters = gstr_filters)
        def fn_validate(key):
            if not self._LsAndGermInfoSet: return []

            #only whack-a-mole plots for the length-1 germs are available
            len1GermFirstEls = [ g[0] for g in self.gatestring_lists['germs']
                                 if len(g) == 1 ]

            keys = ["whack%sMoleBoxes" % gl for gl in len1GermFirstEls]
            if key == expr2: return keys # all computable keys
            elif key in keys: return [key]
            else: return []
        fns[expr2] = (fn, fn_validate)


        expr3 = "whack(.+?)MoleBoxesSummed"
        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            Ls,germs, gsBest, fpr_filters, gstr_filters, _, _, baseStr_dict, strs, st = plot_setup()
            highestL = Ls[-1]; hammerWeight = 10.0; mpc = getMPC()
            gateLabel = _re.match(expr3,key).group(1)
            strToWhack = _gs.GateString( (gateLabel,)*highestL )
            whackAMolePlotFn = getWhackAMolePlotFn()
            return whackAMolePlotFn(strToWhack, self.gatestring_lists['all'],
                                    Ls[st:], germs, baseStr_dict, self.dataset,
                                    gsBest, strs, r"$L$", "germ", scale=1.0,
                                    sumUp=True, title="",whackWith=hammerWeight,
                                    save_to="", minProbClipForWeighting=mpc,
                                    ticSize=20, fidpair_filters=fpr_filters,
                                    gatestring_filters = gstr_filters)

        def fn_validate(key):
            if not self._LsAndGermInfoSet: return []

            #only whack-a-mole plots for the length-1 germs are available
            len1GermFirstEls = [ g[0] for g in self.gatestring_lists['germs']
                                 if len(g) == 1 ]

            keys = ["whack%sMoleBoxesSummed" % gl for gl in len1GermFirstEls]
            if key == expr3: return keys # all computable keys
            elif key in keys: return [key]
            else: return []
        fns[expr3] = (fn, fn_validate)


        expr4 = "bestGateErrGenBoxes(.+)"
        def fn(key, confidenceLevel, vb):
            #cri = self._get_confidence_region(confidenceLevel)
            noConfidenceLevelDependence(confidenceLevel)
            gateLabel = _re.match(expr4,key).group(1)
            gate = self.gatesets['final estimate'].gates[gateLabel]
            targetGate = self.gatesets['target'].gates[gateLabel]
            basisNm   = self.gatesets['final estimate'].get_basis_name()
            basisDims = self.gatesets['final estimate'].get_basis_dimension()
            assert(basisNm == self.gatesets['target'].get_basis_name())
            return _plotting.gate_matrix_errgen_boxplot(
                gate, targetGate, save_to="", mxBasis=basisNm,
                mxBasisDims=basisDims)

        def fn_validate(key):
            if not self._bEssentialResultsSet: return []
            keys = ["bestGateErrGenBoxes%s" % gl
                    for gl in self.gatesets['final estimate'].gates ]
            if key == expr4: return keys # all computable keys
            elif key in keys: return [key]
            else: return []
        fns[expr4] = (fn, fn_validate)


        expr5 = "targetGateBoxes(.+)"
        def fn(key, confidenceLevel, vb):
            #cri = self._get_confidence_region(confidenceLevel)
            noConfidenceLevelDependence(confidenceLevel)
            gateLabel = _re.match(expr5,key).group(1)
            gate = self.gatesets['target'].gates[gateLabel]
            return _plotting.gate_matrix_boxplot(gate, save_to="",
                mxBasis=self.gatesets['target'].get_basis_name(),
                mxBasisDims=self.gatesets['target'].get_basis_dimension())

        def fn_validate(key):
            if not self._bEssentialResultsSet: return []
            keys = ["targetGateBoxes%s" % gl
                    for gl in self.gatesets['final estimate'].gates ]
            if key == expr5: return keys # all computable keys
            elif key in keys: return [key]
            else: return []
        fns[expr5] = (fn, fn_validate)

        expr6 = "bestEstimatePolar(.+?)EvalPlot"
        def fn(key, confidenceLevel, vb):
            #cri = self._get_confidence_region(confidenceLevel)
            noConfidenceLevelDependence(confidenceLevel)
            gateLabel = _re.match(expr6,key).group(1)
            gate = self.gatesets['final estimate'].gates[gateLabel]
            target_gate = self.gatesets['target'].gates[gateLabel]
            return _plotting.polar_eigenval_plot(gate, target_gate,
                                                 title=gateLabel, save_to="")

        def fn_validate(key):
            if not self._bEssentialResultsSet: return []
            keys = ["bestEstimatePolar%sEvalPlot" % gl
                    for gl in self.gatesets['final estimate'].gates ]
            if key == expr6: return keys # all computable keys
            elif key in keys: return [key]
            else: return []
        fns[expr6] = (fn, fn_validate)


        expr7 = "pauliProdHamiltonianDecompBoxes(.+)"
        def fn(key, confidenceLevel, vb):
            #cri = self._get_confidence_region(confidenceLevel)
            noConfidenceLevelDependence(confidenceLevel)
            gateLabel = _re.match(expr7,key).group(1)
            gate = self.gatesets['final estimate'].gates[gateLabel]
            target_gate = self.gatesets['target'].gates[gateLabel]
            basisNm   = self.gatesets['final estimate'].get_basis_name()
            assert(basisNm == self.gatesets['target'].get_basis_name())
            return _plotting.pauliprod_hamiltonian_boxplot(
                gate, target_gate, save_to="", mxBasis=basisNm, boxLabels=True)

        def fn_validate(key):
            if not self._bEssentialResultsSet: return []
            keys = ["pauliProdHamiltonianDecompBoxes%s" % gl
                    for gl in self.gatesets['final estimate'].gates ]
            if key == expr7: return keys # all computable keys
            elif key in keys: return [key]
            else: return []
        fns[expr7] = (fn, fn_validate)

        return fns


    def _get_special_fns(self):
        """
        Return a dictionary of functions which create "special objects"
        identified by the dictionary key.  These functions are used for
        the lazy creation of these objects within the "_specials" member
        of a Results instance.
        """

        def validate_essential(key):
            return [key] if self._bEssentialResultsSet else []
        def validate_LsAndGerms(key):
            return [key] if (self._bEssentialResultsSet and
                             self._LsAndGermInfoSet) else []

        def noConfidenceLevelDependence(level):
            """ Designates a "special" as independent of the confidence level"""
            if level is not None: raise _ResultCache.NoCRDependenceError

        fns = _collections.OrderedDict()


        def fn(key, confidenceLevel, vb):
            printer = VerbosityPrinter.build_printer(vb)
            noConfidenceLevelDependence(confidenceLevel)

            gsTarget = self.gatesets['target']
            gsBestEstimate = self.gatesets['final estimate']

            printer.log("Performing gauge transforms for appendix...")

            best_gs_gauges = _collections.OrderedDict()

            best_gs_gauges['Target'] = _optimizeGauge(
                gsBestEstimate, gsTarget, {'gates': 1.0, 'spam': 1.0},
                verbosity=vb)

            best_gs_gauges['TargetSpam'] = _optimizeGauge(
                gsBestEstimate, gsTarget, {'gates': 1e-3, 'spam': 1.0},
                verbosity=vb)

            best_gs_gauges['TargetGates'] = _optimizeGauge(
                gsBestEstimate, gsTarget, {'gates': 1.0, 'spam': 1e-3},
                verbosity=vb)

            best_gs_gauges['CPTP'] = _optimizeGauge(
                gsBestEstimate, gsTarget, CPpenalty=1e5, TPpenalty=1e5,
                validSpamPenalty=1e5, verbosity=vb)

            best_gs_gauges['TP'] = _optimizeGauge(
                gsBestEstimate, gsTarget, TPpenalty=1e5, verbosity=vb)

            return best_gs_gauges
        fns['gaugeOptAppendixGatesets'] = (fn, validate_essential)


        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)

            best_gs_gauges = self._specials.get('gaugeOptAppendixGatesets',
                                                verbosity=vb)
            gsTarget = self.gatesets['target']

            ret = {}

            for gaugeKey,gopt_gs in best_gs_gauges.items():
                #FUTURE: add confidence region support to these appendices?
                # -- would need to compute confidenceRegionInfo (cri)
                #    for each gauge-optimized gateset, gopt_gs and pass
                #    to appropriate functions below
                ret['best%sGatesetSpamTable' % gaugeKey] = \
                    _generation.get_gateset_spam_table(gopt_gs)
                ret['best%sGatesetSpamParametersTable' % gaugeKey] = \
                    _generation.get_gateset_spam_parameters_table(gopt_gs)
                ret['best%sGatesetGatesTable' % gaugeKey] = \
                    _generation.get_gateset_gates_table(gopt_gs)
                ret['best%sGatesetChoiTable' % gaugeKey] = \
                    _generation.get_gateset_choi_table(gopt_gs)
                ret['best%sGatesetDecompTable' % gaugeKey] = \
                    _generation.get_gateset_decomp_table(gopt_gs)
                ret['best%sGatesetRotnAxisTable' % gaugeKey] = \
                    _generation.get_gateset_rotn_axis_table(gopt_gs)
                ret['best%sGatesetClosestUnitaryTable' % gaugeKey] = \
                    _generation.get_gateset_closest_unitary_table(gopt_gs)
                ret['best%sGatesetVsTargetTable' % gaugeKey] = \
                    _generation.get_gates_vs_target_table(
                    gopt_gs, gsTarget, None)
                ret['best%sGatesetErrorGenTable' % gaugeKey] = \
                    _generation.get_gates_vs_target_err_gen_table(
                    gopt_gs, gsTarget, self.options.errgen_type)

            return ret
        fns['gaugeOptAppendixTables'] = (fn, validate_essential)


        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)
            ret = {}

            for gaugeKey in ('Target','TargetSpam','TargetGates','CPTP','TP'):
                ret['best%sGatesetSpamTable' % gaugeKey] = \
                    _generation.get_blank_table()
                ret['best%sGatesetSpamParametersTable' % gaugeKey] = \
                    _generation.get_blank_table()
                ret['best%sGatesetGatesTable' % gaugeKey] = \
                    _generation.get_blank_table()
                ret['best%sGatesetChoiTable' % gaugeKey] = \
                    _generation.get_blank_table()
                ret['best%sGatesetDecompTable' % gaugeKey] = \
                    _generation.get_blank_table()
                ret['best%sGatesetRotnAxisTable' % gaugeKey] = \
                    _generation.get_blank_table()
                ret['best%sGatesetClosestUnitaryTable' % gaugeKey] = \
                    _generation.get_blank_table()
                ret['best%sGatesetVsTargetTable' % gaugeKey] = \
                    _generation.get_blank_table()
                ret['best%sGatesetErrorGenTable' % gaugeKey] = \
                    _generation.get_blank_table()

            return ret
        fns['blankGaugeOptAppendixTables'] = (fn, validate_essential)


        def fn(key, confidenceLevel, vb):
            printer = VerbosityPrinter.build_printer(vb)
            noConfidenceLevelDependence(confidenceLevel)

            gsTarget = self.gatesets['target']
            gsBestEstimate = self.gatesets['final estimate']

            best_gs_gauges = _collections.OrderedDict()

            for gateLabel in gsBestEstimate.gates:
                best_gs_gauges[gateLabel] = _optimizeGauge(
                    gsBestEstimate, gsTarget, 
                    {'gates': 0.0, 'spam': 0.0, gateLabel: 1.0},
                    verbosity=vb)

            return best_gs_gauges
        fns['singleGateTargetGaugeOptGatesets'] = (fn, validate_essential)


        def fn(key, confidenceLevel, vb):
            printer = VerbosityPrinter.build_printer(vb)
            noConfidenceLevelDependence(confidenceLevel)

            gsTarget = self.gatesets['target']
            gsBestEstimate = self.gatesets['final estimate']

            #Heusistic parameters for CPTP gauge opt that doesn't take too long
            if hasattr(self.parameters['gaugeOptParams'],"keys"):
                gaugeParams = self.parameters['gaugeOptParams'].copy()
            else:
                gaugeParams = self.parameters['gaugeOptParams'][0].copy()
            gaugeParams['CPpenalty'] = 100
            gaugeParams['TPpenalty'] = 100
            gaugeParams['validSpamPenalty'] = 0
            gaugeParams['tol'] = 0.1
            gaugeParams['maxiter'] = 100
            gaugeParams['method'] = 'BFGS'
            gaugeParams['targetGateset'] = gsTarget
            #gaugeParams['verbosity'] = 5 #DEBUG
            return _optimizeGauge(gsBestEstimate, **gaugeParams)
        fns['CPTPGaugeOptGateset'] = (fn, validate_essential)


        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)

            direct_specs = _ssc.build_spam_specs(
                prepStrs=self.gatestring_lists['prep fiducials'],
                effectStrs=self.gatestring_lists['effect fiducials'],
                prep_labels=self.gatesets['target'].get_prep_labels(),
                effect_labels=self.gatesets['target'].get_effect_labels() )

            baseStrs = [] # (L,germ) base strings without duplicates
            fullDict = self.parameters['L,germ tuple base string dict']
            for L in self.parameters['max length list']:
                for germ in self.gatestring_lists['germs']:
                    if fullDict[(L,germ)] not in baseStrs:
                        baseStrs.append( fullDict[(L,germ)] )

            return _plotting.direct_lgst_gatesets(
                baseStrs, self.dataset, direct_specs, self.gatesets['target'],
                svdTruncateTo=4, verbosity=0)
                #TODO: svdTruncateTo set elegantly?
        fns["direct_lgst_gatesets"] = (fn, validate_LsAndGerms)


        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)

            gsTarget = self.gatesets['target']
            direct_specs = _ssc.build_spam_specs(
                prepStrs=self.gatestring_lists['prep fiducials'],
                effectStrs=self.gatestring_lists['effect fiducials'],
                prep_labels=gsTarget.get_prep_labels(),
                effect_labels=gsTarget.get_effect_labels() )

            baseStrs = [] # (L,germ) base strings without duplicates
            fullDict = self.parameters['L,germ tuple base string dict']
            for L in self.parameters['max length list']:
                for germ in self.gatestring_lists['germs']:
                    if fullDict[(L,germ)] not in baseStrs:
                        baseStrs.append( fullDict[(L,germ)] )

            if self.parameters['objective'] == "chi2":
                mpc = self.parameters['minProbClipForWeighting']
                return _plotting.direct_mc2gst_gatesets(
                    baseStrs, self.dataset, direct_specs, gsTarget,
                    svdTruncateTo=gsTarget.get_dimension(),
                    minProbClipForWeighting=mpc,
                    probClipInterval=self.parameters['probClipInterval'],
                    verbosity=0)

            elif self.parameters['objective'] == "logl":
                mpc = self.parameters['minProbClip']
                return _plotting.direct_mlgst_gatesets(
                    baseStrs, self.dataset, direct_specs, gsTarget,
                    svdTruncateTo=gsTarget.get_dimension(),
                    minProbClip=mpc,
                    probClipInterval=self.parameters['probClipInterval'],
                    verbosity=0)
            else:
                raise ValueError("Invalid Objective: %s" %
                                 self.parameters['objective'])
        fns["DirectLongSeqGatesets"] = (fn, validate_LsAndGerms)


        def fn(key, confidenceLevel, vb):
            noConfidenceLevelDependence(confidenceLevel)

            baseStr_dict = self._getBaseStrDict()
            strs  = (self.gatestring_lists['prep fiducials'],
                     self.gatestring_lists['effect fiducials'])
            germs = self.gatestring_lists['germs']
            gsBest = self.gatesets['final estimate']
            fidPairs = self.parameters['fiducial pairs']
            Ls = self.parameters['max length list']
            st = 1 if Ls[0] == 0 else 0 #start index: skip LGST column in plots

            if fidPairs is None: fidpair_filters = None
            elif isinstance(fidPairs,dict) or hasattr(fidPairs,"keys"):
                #Assume fidPairs is a dict indexed by germ
                fidpair_filters = { (x,y): fidPairs[germ] 
                                    for x in Ls[st:] for y in germs }
            else:
                #Assume fidPairs is a list
                fidpair_filters = { (x,y): fidPairs
                                    for x in Ls[st:] for y in germs }

            if fidPairs is None: fidpair_filters = None
            gstr_filters = { (x,y) : self.gatestring_lists['iteration'][i]
                             for i,x in enumerate(Ls[st:],start=st)
                             for y in germs }

            obj = self.parameters['objective']
            assert(obj in ("chi2","logl"))
            if obj == "chi2":
                plotFn = _plotting.chi2_boxplot
                mpc = self.parameters['minProbClipForWeighting']

            elif obj == "logl":
                plotFn = _plotting.logl_boxplot
                mpc = self.parameters['minProbClip']

            maxH = 9.0 # max inches for graphic
            minboxH = 0.075 #min height per box
            germH = (len(self.gatestring_lists['effect fiducials'])+1)*minboxH # + 1 for space
            maxGermsPerFig = max(int(maxH / germH - 2), 1)
            figs = []; n = 0
            while( n < len(germs) ):
                fig_germs = list(reversed(germs[n:n+maxGermsPerFig]))
                fig = plotFn(Ls[st:], fig_germs, baseStr_dict,
                             self.dataset, gsBest, strs,
                             r"$L$", "germ", scale=1.0, sumUp=False,
                             histogram=True, title="", 
                             fidpair_filters=fidpair_filters,
                             gatestring_filters = gstr_filters,
                             linlg_pcntle=float(self.parameters['linlogPercentile']) / 100,
                             minProbClipForWeighting=mpc, save_to="", ticSize=20)
                figs.append(fig); n += maxGermsPerFig

            return figs

        fns["bestEstimateColorBoxPlotPages"] = (fn,validate_LsAndGerms)


        return fns


#    def set_additional_info(self,minProbClip=1e-6, minProbClipForWeighting=1e-4,
#                          probClipInterval=(-1e6,1e6), radius=1e-4,
#                          weightsDict=None, defaultDirectory=None, defaultBasename=None,
#                          mxBasis="gm"):
#        """
#        Set advanced parameters for producing derived outputs.  Usually the default
#        values are fine (which is why setting these inputs is separated into a
#        separate function).
#
#        Parameters
#        ----------
#        minProbClip : float, optional
#            The minimum probability treated normally in the evaluation of the log-likelihood.
#            A penalty function replaces the true log-likelihood for probabilities that lie
#            below this threshold so that the log-likelihood never becomes undefined (which improves
#            optimizer performance).
#
#        minProbClipForWeighting : float, optional
#            Sets the minimum and maximum probability p allowed in the chi^2 weights: N/(p*(1-p))
#            by clipping probability p values to lie within the interval
#            [ minProbClipForWeighting, 1-minProbClipForWeighting ].
#
#        probClipInterval : 2-tuple or None, optional
#           (min,max) values used to clip the probabilities predicted by gatesets during the
#           least squares search for an optimal gateset (if not None).
#
#        radius : float, optional
#           Specifies the severity of rounding used to "patch" the zero-frequency
#           terms of the log-likelihood.
#
#        weightsDict : dict, optional
#           A dictionary with keys == gate strings and values == multiplicative scaling
#           factor for the corresponding gate string. The default is no weight scaling at all.
#
#        defaultDirectory : string, optional
#           Path to the default directory for generated reports and presentations.
#
#        defaultBasename : string, optional
#           Default basename for generated reports and presentations.
#
#        mxBasis : {"std", "gm", "pp"}, optional
#           The basis used to interpret the gate matrices.
#
#
#        Returns
#        -------
#        None
#        """
#
#
#        self.additionalInfo = { 'weights': weightsDict,
#                                'minProbClip': minProbClip,
#                                'minProbClipForWeighting': minProbClipForWeighting,
#                                'probClipInterval': probClipInterval,
#                                'radius': radius,
#                                'hessianProjection': 'std',
#                                'defaultDirectory': defaultDirectory,
#                                'defaultBasename': defaultBasename,
#                                'mxBasis': "gm"}
#
#    def set_template_path(self, pathToTemplates):
#        """
#        Sets the location of GST report and presentation templates.
#
#        Parameters
#        ----------
#        pathToTemplates : string
#           The path to a folder containing GST's template files.
#           Usually this can be determined automatically (the default).
#        """
#        self.options.template_path = pathToTemplates
#
#
#    def set_latex_cmd(self, latexCmd):
#        """
#        Sets the shell command used for compiling latex reports and
#        presentations.
#
#        Parameters
#        ----------
#        latexCmd : string
#           The command to run to invoke the latex compiler,
#           typically just 'pdflatex' when it is on the system
#           path.
#        """
#        self.latexCmd = latexCmd


    def _get_confidence_region(self, confidenceLevel):
        """
        Get the ConfidenceRegion object associated with a given confidence level.
        One will be created and cached the first time given level is requested,
        and future requests will return the cached object.

        Parameters
        ----------
        confidenceLevel : float
           The confidence level (between 0 and 100).

        Returns
        -------
        ConfidenceRegion
        """

        assert(self._bEssentialResultsSet)

        if confidenceLevel is None:
            return None

        if confidenceLevel not in self._confidence_regions:

            #Negative confidence levels ==> non-Markovian error bars
            if confidenceLevel < 0:
                actual_confidenceLevel = -confidenceLevel
                regionType = "non-markovian"
            else:
                actual_confidenceLevel = confidenceLevel
                regionType = "std"

            if self.parameters['objective'] == "logl":
                cr = _generation.get_logl_confidence_region(
                    self.gatesets['final estimate'], self.dataset,
                    actual_confidenceLevel,
                    self.gatestring_lists['final'],
                    self.parameters['probClipInterval'],
                    self.parameters['minProbClip'],
                    self.parameters['radius'],
                    self.parameters['hessianProjection'],
                    regionType, self._comm,
                    self.parameters['memLimit'],
                    self.parameters['cptpPenaltyFactor'],
                    self.parameters['distributeMethod'])
            elif self.parameters['objective'] == "chi2":
                cr = _generation.get_chi2_confidence_region(
                    self.gatesets['final estimate'], self.dataset,
                    actual_confidenceLevel,
                    self.gatestring_lists['final'],
                    self.parameters['probClipInterval'],
                    self.parameters['minProbClipForWeighting'],
                    self.parameters['hessianProjection'],
                    regionType, self._comm,
                    self.parameters['memLimit'])
            else:
                raise ValueError("Invalid objective given in essential" +
                                 " info: %s" % self.parameters['objective'])
            self._confidence_regions[confidenceLevel] = cr

        return self._confidence_regions[confidenceLevel]


    def _merge_template(self, qtys, templateFilename, outputFilename):
        if self.options.template_path is None:
            templateFilename = _os.path.join( _os.path.dirname(_os.path.abspath(__file__)),
                                              "templates", templateFilename )
        else:
            templateFilename = _os.path.join( self.options.template_path,
                                              templateFilename )

        template = ''
        with open(templateFilename, 'r') as templatefile:
            template = templatefile.read()
        template = template.replace("{", "{{").replace("}", "}}") #double curly braces (for format processing)

        # Replace template field markers with `str.format` fields.
        template = _re.sub( r"\\putfield\{\{([^}]+)\}\}\{\{[^}]*\}\}", "{\\1}", template)

        # Replace str.format fields with values and write to output file
        template = template.format(**qtys)
        with open(outputFilename, 'w') as outputfile:
            outputfile.write(template)

    def _getBaseStrDict(self, remove_dups = True):
        #if remove_dups == True, remove duplicates in
        #  L_germ_tuple_to_baseStr_dict by replacing with None

        assert(self._bEssentialResultsSet)
        assert(self._LsAndGermInfoSet)

        baseStr_dict = _collections.OrderedDict()
        st = 1 if self.parameters['max length list'][0] == 0 else 0
          #start index: skips LGST column in report color box plots

        tmpRunningList = []
        fullDict = self.parameters['L,germ tuple base string dict']
        for L in self.parameters['max length list'][st:]:
            for germ in self.gatestring_lists['germs']:
                if remove_dups and fullDict[(L,germ)] in tmpRunningList:
                    baseStr_dict[(L,germ)] = None
                else:
                    tmpRunningList.append( fullDict[(L,germ)] )
                    baseStr_dict[(L,germ)] = fullDict[(L,germ)]
        return baseStr_dict




    def create_full_report_pdf(self, confidenceLevel=None, filename="auto",
                            title="auto", datasetLabel="auto", suffix="",
                            debugAidsAppendix=False, gaugeOptAppendix=False,
                            pixelPlotAppendix=False, whackamoleAppendix=False,
                            m=0, M=10, tips=False, verbosity=0, comm=None):
        """
        Create a "full" GST report.  This report is the most detailed of any of
        the GST reports, and includes background and explanation text to help
        the user interpret the contained results.

        Parameters
        ----------
        confidenceLevel : float, optional
           If not None, then the confidence level (between 0 and 100) used in
           the computation of confidence regions/intervals. If None, no
           confidence regions or intervals are computed.

        filename : string, optional
           The output filename where the report file(s) will be saved.  Specifying
           "auto" will use the default directory and base name (specified in
           set_additional_info) if given, otherwise the file "GSTReport.pdf" will
           be output to the current directoy.

        title : string, optional
           The title of the report.  "auto" uses a default title which
           specifyies the label of the dataset as well.

        datasetLabel : string, optional
           A label given to the dataset.  If set to "auto", then the label
           will be the base name of the dataset filename without extension
           (if given) or "$\\mathcal{D}$" (if not).

        suffix : string, optional
           A suffix to add to the end of the report filename.  Useful when
           filename is "auto" and you generate different reports using
           the same dataset.

        debugAidsAppendix : bool, optional
           Whether to include the "debugging aids" appendix.  This
           appendix contains comparisons of GST and Direct-GST and small-
           eigenvalue error rates among other quantities potentially
           useful for figuring out why the GST estimate did not fit
           the data as well as expected.

        gaugeOptAppendix : bool, optional
           Whether to include the "gauge optimization" appendix.  This
           appendix shows the results of gauge optimizing GST's best
           estimate gate set in several different ways, and thus shows
           how various report quantities can vary by a different gauge
           choice.

        pixelPlotAppendix : bool, optional
           Whether to include the "pixel plots" appendix, which shows
           the goodness of fit, in the form of color box plots, for the
           intermediate iterations of the GST algortihm.

        whackamoleAppendix : bool, optional
           Whether to include the "whack-a-mole" appendix, which contains
           colr box plots showing the effect of reducing ("whacking") one
           particular part of the overall goodness of fit box plot.

        m, M : float, optional
           Minimum and Maximum values of the color scale used in the report's
           color box plots.

        tips : boolean, optional
           If True, additional markup and tooltips are included in the produced
           PDF which indicate how tables and figures in the report correspond
           to members of this Results object.  These "tips" can be useful if
           you want to further manipulate the data contained in a table or
           figure.

        verbosity : int, optional
           How much detail to send to stdout.

        comm : mpi4py.MPI.Comm, optional
            When not None, an MPI communicator for distributing the computation
            across multiple processors.

        Returns
        -------
        None
        """
        printer = VerbosityPrinter.build_printer(verbosity, comm=comm)

        assert(self._bEssentialResultsSet)
        self.confidence_level = confidenceLevel
        self._comm = comm
          #set "current" level, used by ResultCache member dictionaries

        if tips:
            def tooltiptex(directive):
                return '\\pdftooltip{{\\color{blue}\\texttt{%s}}\\quad}' \
                    % directive + '{Access this information in pyGSTi via' \
                    + '<ResultsObject>%s}' % directive


        else:
            def tooltiptex(directive):
                return "" #tooltips disabled


        #Get report output filename
        default_dir = self.parameters['defaultDirectory']
        default_base = self.parameters['defaultBasename']

        if filename != "auto":
            report_dir = _os.path.dirname(filename)
            report_base = _os.path.splitext( _os.path.basename(filename) )[0] \
                           + suffix
        else:
            cwd = _os.getcwd()
            report_dir  = default_dir  if (default_dir  is not None) else cwd
            report_base = default_base if (default_base is not None) \
                                       else "GSTReport"
            report_base += suffix

        if datasetLabel == "auto":
            if default_base is not None:
                datasetLabel = _latex.latex_escaped( default_base )
            else:
                datasetLabel = "$\\mathcal{D}$"

        if title == "auto": title = "GST report for %s" % datasetLabel

        ######  Generate Report ######
        #Steps:
        # 1) generate latex tables
        # 2) generate plots
        # 3) populate template latex file => report latex file
        # 4) compile report latex file into PDF
        # 5) remove auxiliary files generated during compilation
        #  FUTURE?? determine what we need to compute & plot by reading
        #           through the template file?

        #Note: for now, we assume the best gateset corresponds to the last
        #      L-value
        # best_gs = self.gatesets['final estimate']
        # v = verbosity # shorthand

        if not self._LsAndGermInfoSet: #cannot create appendices
            debugAidsAppendix = False  # which depend on this structure
            pixelPlotAppendix = False
            whackamoleAppendix = False

        qtys = {} # dictionary to store all latex strings
                  # to be inserted into report template
        qtys['title'] = title
        qtys['datasetLabel'] = datasetLabel
        qtys['settoggles'] =  "\\toggle%s{confidences}\n" % \
            ("false" if confidenceLevel is None else "true")
        qtys['settoggles'] += "\\toggle%s{LsAndGermsSet}\n" % \
            ("true" if self._LsAndGermInfoSet else "false")
        qtys['settoggles'] += "\\toggle%s{debuggingaidsappendix}\n" % \
            ("true" if debugAidsAppendix else "false")
        qtys['settoggles'] += "\\toggle%s{gaugeoptappendix}\n" % \
            ("true" if gaugeOptAppendix else "false")
        qtys['settoggles'] += "\\toggle%s{pixelplotsappendix}\n" % \
            ("true" if pixelPlotAppendix else "false")
        qtys['settoggles'] += "\\toggle%s{whackamoleappendix}\n" % \
            ("true" if whackamoleAppendix else "false")
        qtys['confidenceLevel'] = "%g" % \
            confidenceLevel if confidenceLevel is not None else "NOT-SET"
        qtys['linlg_pcntle'] = self.parameters['linlogPercentile']

        if self.options.errgen_type == "logTiG":
            qtys['errorgenformula'] = "$\hat{G} = G_{\mathrm{target}}e^{\mathbb{L}}$"
        elif self.options.errgen_type == "logG-logT":
            qtys['errorgenformula'] = "$\hat{G} = e^{\mathbb{L} + \log G_{\mathrm{target}}}$"
        else:
            qtys['errorgenformula'] = "???"

        if confidenceLevel is not None:
            cri = self._get_confidence_region(confidenceLevel)
            qtys['confidenceIntervalScaleFctr'] = \
                "%.3g" % cri.intervalScaling
            qtys['confidenceIntervalNumNonGaugeParams'] = \
                "%d" % cri.nNonGaugeParams
        else:
            cri = None
            qtys['confidenceIntervalScaleFctr'] = "NOT-SET"
            qtys['confidenceIntervalNumNonGaugeParams'] = "NOT-SET"

        pdfInfo = [('Author','pyGSTi'), ('Title', title),
                   ('Keywords', 'GST'), ('pyGSTi Version',_version.__version__),
                   ('opt_long_tables', self.options.long_tables),
                   ('opt_table_class', self.options.table_class),
                   ('opt_template_path', self.options.template_path),
                   ('opt_latex_cmd', self.options.latex_cmd) ]
                   #('opt_latex_postcmd', self.options.latex_postcmd) #TODO: add this

        # Note: use definite ordering of parameters
        for key in sorted(list(self.parameters.keys())):
            pdfInfo.append( (key, self.parameters[key]) )
        qtys['pdfinfo'] = _to_pdfinfo( pdfInfo )


        #Get figure directory for figure generation *and* as a
        # scratch space for tables.
        D = report_base + "_files" #figure directory relative to reportDir
        if not _os.path.isdir( _os.path.join(report_dir,D)):
            _os.mkdir( _os.path.join(report_dir,D))


        # 1) get latex tables
        printer.log("*** Generating tables ***")

        std_tables = \
            ('targetSpamTable','targetGatesTable','datasetOverviewTable',
             'bestGatesetSpamTable','bestGatesetSpamParametersTable',
             'bestGatesetGaugeOptParamsTable',
             'bestGatesetGatesTable','bestGatesetChoiTable',
             'bestGatesetDecompTable','bestGatesetRotnAxisTable',
             'bestGatesetVsTargetTable','bestGatesetErrorGenTable')
             #removed: 'bestGatesetClosestUnitaryTable',

        ls_and_germs_tables = ('fiducialListTable','prepStrListTable',
                               'effectStrListTable','germListTable',
                               'progressTable')

        tables_to_compute = std_tables
        tables_to_blank = []

        if self._LsAndGermInfoSet:
            tables_to_compute += ls_and_germs_tables
        else:
            tables_to_blank += ls_and_germs_tables

        with printer.progress_logging(1):
            for i, key in enumerate(tables_to_compute):
                printer.show_progress(i, len(tables_to_compute), prefix='', end='')
                qtys[key] = self.tables.get(key, verbosity=printer - 1).render(
                    'latex',longtables=self.options.long_tables, scratchDir=D,
                    precision=self.options.precision,
                    polarprecision=self.options.polar_precision,
                    sciprecision=self.options.sci_precision)
                qtys["tt_"+key] = tooltiptex(".tables['%s']" % key)

        for key in tables_to_blank:
            qtys[key] = _generation.get_blank_table().render(
                'latex',longtables=self.options.long_tables)
            qtys["tt_"+key] = ""

        #get appendix tables if needed
        if gaugeOptAppendix:
            goaTables = self._specials.get('gaugeOptAppendixTables',verbosity=printer - 1)
            qtys.update( { key : goaTables[key].render(
                        'latex', longtables=self.options.long_tables,
                        scratchDir=D,
                        precision=self.options.precision,
                        polarprecision=self.options.polar_precision,
                        sciprecision=self.options.sci_precision)
                           for key in goaTables }  )
            #TODO: tables[ref] and then tooltips?

        elif any((debugAidsAppendix, pixelPlotAppendix, whackamoleAppendix)):
            goaTables = self._specials.get('blankGaugeOptAppendixTables',
                              verbosity=printer - 1)   # fill keys with blank tables
            qtys.update( { key : goaTables[key].render(
                        'latex',longtables=self.options.long_tables,
                        precision=self.options.precision,
                        polarprecision=self.options.polar_precision,
                        sciprecision=self.options.sci_precision)
                           for key in goaTables }  )  # for format substitution
            #TODO: tables[ref] and then tooltips?


        # 2) generate plots
        printer.log("*** Generating plots ***")

        if _matplotlib.is_interactive():
            _matplotlib.pyplot.ioff()
            bWasInteractive = True
        else: bWasInteractive = False

        maxW,maxH = 6.5,8.0 #max width and height of graphic in latex document (in inches)

        def incgr(figFilenm,W=None,H=None): #includegraphics "macro"
            if W is None: W = maxW
            if H is None: H = maxH
            return "\\includegraphics[width=%.2fin,height=%.2fin" % (W,H) + \
                ",keepaspectratio]{%s/%s}" % (D,figFilenm)

        def set_fig_qtys(figkey, figFilenm, v, W=None,H=None):
            fig = self.figures.get(figkey, verbosity=v)
            fig.save_to(_os.path.join(report_dir, D, figFilenm))
            qtys[figkey] = incgr(figFilenm,W,H)
            qtys['tt_' + figkey] = tooltiptex(".figures['%s']" % figkey)
            return fig

        #Chi2 or logl plots
        if self._LsAndGermInfoSet:
            Ls = self.parameters['max length list']
            st = 1 if Ls[0] == 0 else 0 #start index: skip LGST column in plots
            nPlots = (len(Ls[st:])-1)+2 if pixelPlotAppendix else 2

            if self.parameters['objective'] == "chi2":
                plotFnName,plotFnLatex = "Chi2", "$\chi^2$"
            elif self.parameters['objective'] == "logl":
                plotFnName,plotFnLatex = "LogL", "$\\log(\\mathcal{L})$"
            else:
                raise ValueError("Invalid objective value: %s"
                                 % self.parameters['objective'])
            printer.log("%s plots (%d): " % (plotFnName, nPlots))

            with printer.progress_logging(1):
                printer.show_progress(0, 3, prefix='', end='')

                w = min(len(self.gatestring_lists['prep fiducials']) * 0.3,maxW)
                h = min(len(self.gatestring_lists['effect fiducials']) * 0.3,maxH)
                fig = set_fig_qtys("colorBoxPlotKeyPlot",
                                   "colorBoxPlotKey.png", printer - 1, w,h)

                printer.show_progress(1, 3, prefix='', end='')

                fig = set_fig_qtys("bestEstimateColorBoxPlot",
                                   "best%sBoxes.pdf" % plotFnName, printer - 1)
                maxX = fig.get_extra_info()['nUsedXs']
                maxY = fig.get_extra_info()['nUsedYs']

                #qtys["bestEstimateColorBoxPlot_hist"] = \
                #    incgr("best%sBoxes_hist.pdf" % plotFnName figFilenm)
                #    #no tooltip for histogram... - probably should make it
                #    # it's own element of .figures dict

                printer.show_progress(2, 3, prefix='', end='')
                fig = set_fig_qtys("invertedBestEstimateColorBoxPlot",
                                   "best%sBoxes_inverted.pdf" % plotFnName, printer - 1)
        else:
            for figkey in ["colorBoxPlotKeyPlot",
                           "bestEstimateColorBoxPlot",
                           "invertedBestEstimateColorBoxPlot"]:
                qtys[figkey] = qtys["tt_"+figkey] = ""


        pixplots = ""
        if pixelPlotAppendix:
            Ls = self.parameters['max length list']
            with printer.progress_logging(1):
                for i in range(st,len(Ls)-1):

                    printer.show_progress(i, len(Ls)-1, prefix='', end='') # -2 is intentional

                    fig = self.figures.get("estimateForLIndex%dColorBoxPlot" % i,
                                           verbosity=printer)
                    fig.save_to( _os.path.join(report_dir, D,
                                               "L%d_%sBoxes.pdf" % (i,plotFnName)))
                    lx = fig.get_extra_info()['nUsedXs']
                    ly = fig.get_extra_info()['nUsedYs']

                    #scale figure size according to number of rows and columns+1
                    # (+1 for labels ~ another col) relative to initial plot
                    W = float(lx+1)/float(maxX+1) * maxW
                    H = float(ly)  /float(maxY)   * maxH

                    pixplots += "\n"
                    pixplots += "\\begin{figure}\n"
                    pixplots += "\\begin{center}\n"
                    pixplots += "\\includegraphics[width=%.2fin,height=%.2fin," \
                        % (W,H) + "keepaspectratio]{%s/L%d_%sBoxes.pdf}\n" \
                        %(D,i,plotFnName)
                    pixplots += \
                        "\\caption{Box plot of iteration %d (L=%d) " % (i,Ls[i]) \
                        + "gateset %s values.\label{L%dGateset%sBoxPlot}}\n" \
                        % (plotFnLatex,i,plotFnName)
                    #TODO: add conditional tooltip string to start of caption
                    pixplots += "\\end{center}\n"
                    pixplots += "\\end{figure}\n"

        #Set template quantity (empty string if appendix disabled)
        qtys['intermediate_pixel_plot_figures'] = pixplots

        printer.log("")

        if debugAidsAppendix:
            #DirectLGST and deviation
            printer.log(" -- Direct-X plots ", end='')
            printer.log("(2):")

            #if verbosity > 0:
            #    print " ?",; _sys.stdout.flush()
            #fig = set_fig_qtys("directLGSTColorBoxPlot",
            #                   "directLGST%sBoxes.pdf" % plotFnName)

            with printer.progress_logging(1):

            #if verbosity > 0:
            #    print " ?",; _sys.stdout.flush()
            #fig = set_fig_qtys("directLGSTDeviationColorBoxPlot",
            #                   "directLGSTDeviationBoxes.pdf",W=4,H=5)
                printer.show_progress(0, 2, prefix='', end='')
                fig = set_fig_qtys("directLongSeqGSTColorBoxPlot",
                               "directLongSeqGST%sBoxes.pdf" % plotFnName, printer - 1)

                #if verbosity > 0:
                #    print " ?",; _sys.stdout.flush()
                #fig = set_fig_qtys("directLGSTDeviationColorBoxPlot",
                #                   "directLGSTDeviationBoxes.pdf",W=4,H=5)

                printer.show_progress(1, 2, prefix='', end='')
                fig = set_fig_qtys("directLongSeqGSTDeviationColorBoxPlot",
                                   "directLongSeqGSTDeviationBoxes.pdf", printer - 1, W=4,H=5)

                printer.log('')

            #Small eigenvalue error rate
            printer.log(" -- Error rate plots...")
            fig = set_fig_qtys("smallEigvalErrRateColorBoxPlot",
                               "smallEigvalErrRateBoxes.pdf", printer - 1, W=4,H=5)
        else:
            #UNUSED: "directLGSTColorBoxPlot", "directLGSTDeviationColorBoxPlot"
            for figkey in ["directLongSeqGSTColorBoxPlot",
                           "directLongSeqGSTDeviationColorBoxPlot",
                           "smallEigvalErrRateColorBoxPlot"]:
                qtys[figkey] = qtys["tt_"+figkey] = ""


        whackamoleplots = ""
        if whackamoleAppendix:
            #Whack-a-mole plots for highest L of each length-1 germ
            Ls = self.parameters['max length list']
            highestL = Ls[-1]; # allGateStrings = self.gatestring_lists['all']
            hammerWeight = 10.0
            len1Germs = [ g for g in self.gatestring_lists['germs']
                          if len(g) == 1 ]

            printer.log(" -- Whack-a-mole plots (%d): " % (2*len(len1Germs)), end='')

            with printer.progress_logging(1):
                for i, germ in enumerate(len1Germs):

                    printer.show_progress(i,  len(len1Germs), prefix='', end='')

                    fig = self.figures.get("whack%sMoleBoxes" % germ[0],verbosity=printer - 1)
                    fig.save_to(_os.path.join(report_dir, D,"whack%sMoleBoxes.pdf"
                                              % germ[0]))

                    whackamoleplots += "\n"
                    whackamoleplots += "\\begin{figure}\n"
                    whackamoleplots += "\\begin{center}\n"
                    whackamoleplots += "\\includegraphics[width=%.2fin,height=%.2fin,keepaspectratio]{%s/whack%sMoleBoxes.pdf}\n" % (maxW,maxH,D,germ[0])
                    whackamoleplots += "\\caption{Whack-a-%s-mole box plot for $\mathrm{%s}^{%d}$." % (plotFnLatex,germ[0],highestL)
                    #TODO: add conditional tooltip string to start of caption
                    whackamoleplots += "  Hitting with hammer of weight %.1f.\label{Whack%sMoleBoxPlot}}\n" % (hammerWeight,germ[0])
                    whackamoleplots += "\\end{center}\n"
                    whackamoleplots += "\\end{figure}\n"

            with printer.progress_logging(1):
                for i,germ in enumerate(len1Germs):
                    printer.show_progress(i, len(len1Germs), prefix='', end='')

                    fig = self.figures.get("whack%sMoleBoxesSummed" % germ[0],
                                           verbosity=printer - 1)
                    fig.save_to(_os.path.join(report_dir, D,"whack%sMoleBoxesSummed.pdf" % germ[0]))

                    whackamoleplots += "\n"
                    whackamoleplots += "\\begin{figure}\n"
                    whackamoleplots += "\\begin{center}\n"
                    whackamoleplots += "\\includegraphics[width=4in,height=5in,keepaspectratio]{%s/whack%sMoleBoxesSummed.pdf}\n" % (D,germ[0])
                    whackamoleplots += "\\caption{Whack-a-%s-mole box plot for $\mathrm{%s}^{%d}$, summed over fiducial matrix." % (plotFnLatex,germ[0],highestL)
                    #TODO: add conditional tooltip string to start of caption
                    whackamoleplots += "  Hitting with hammer of weight %.1f.\label{Whack%sMoleBoxPlotSummed}}\n" % (hammerWeight,germ[0])
                    whackamoleplots += "\\end{center}\n"
                    whackamoleplots += "\\end{figure}\n"

            printer.log('')
        #Set template quantity (empty string if appendix disabled)
        qtys['whackamole_plot_figures'] = whackamoleplots

        if bWasInteractive:
            _matplotlib.pyplot.ion()


        # 3) populate template latex file => report latex file
        printer.log("*** Merging into template file ***")

        mainTexFilename = _os.path.join(report_dir, report_base + ".tex")
        appendicesTexFilename = _os.path.join(report_dir, report_base + "_appendices.tex")
        # TODO: pdffilename is never used
        pdfFilename = _os.path.join(report_dir, report_base + ".pdf")

        if self.parameters['objective'] == "chi2":
            mainTemplate = "report_chi2_main.tex"
            appendicesTemplate = "report_chi2_appendices.tex"
        elif self.parameters['objective'] == "logl":
            mainTemplate = "report_logL_main.tex"
            appendicesTemplate = "report_logL_appendices.tex"
        else:
            raise ValueError("Invalid objective value: %s"
                             % self.parameters['objective'])

        if any( (debugAidsAppendix, gaugeOptAppendix,
                 pixelPlotAppendix, whackamoleAppendix) ):
            qtys['appendices'] = "\\input{%s}" % \
                _os.path.basename(appendicesTexFilename)
            self._merge_template(qtys, appendicesTemplate,
                                 appendicesTexFilename)
        else: qtys['appendices'] = ""
        self._merge_template(qtys, mainTemplate, mainTexFilename)


        # 4) compile report latex file into PDF
        printer.log("Latex file(s) successfully generated.  Attempting to compile with pdflatex...")
        cwd = _os.getcwd()
        if len(report_dir) > 0:
            _os.chdir(report_dir)

        try:
            self._compile_latex_report(report_dir, report_base,
                                       self.options.latex_call, printer)
        except _subprocess.CalledProcessError as e:
            printer.error("pdflatex returned code %d " % e.returncode +
                          "Check %s.log to see details." % report_base)
        finally:
            _os.chdir(cwd)

        return


    def create_brief_report_pdf(self, confidenceLevel=None,
                           filename="auto", title="auto", datasetLabel="auto",
                           suffix="", m=0, M=10, tips=False, verbosity=0,
                           comm=None):
        """
        Create a "brief" GST report.  This report is collects what are typically
        the most relevant tables and plots from the "full" report and presents
        them in an order or efficient analysis by a user familiar with the GST
        full report.  Descriptive text has been all but removed.

        Parameters
        ----------
        confidenceLevel : float, optional
           If not None, then the confidence level (between 0 and 100) used in
           the computation of confidence regions/intervals. If None, no
           confidence regions or intervals are computed.

        filename : string, optional
           The output filename where the report file(s) will be saved.  Specifying
           "auto" will use the default directory and base name (specified in
           set_additional_info) if given, otherwise the file "GSTBrief.pdf" will
           be output to the current directoy.

        title : string, optional
           The title of the report.  "auto" uses a default title which
           specifyies the label of the dataset as well.

        datasetLabel : string, optional
           A label given to the dataset.  If set to "auto", then the label
           will be the base name of the dataset filename without extension
           (if given) or "$\\mathcal{D}$" (if not).

        suffix : string, optional
           A suffix to add to the end of the report filename.  Useful when
           filename is "auto" and you generate different reports using
           the same dataset.

        m, M : float, optional
           Minimum and Maximum values of the color scale used in the report's
           color box plots.

        tips : boolean, optional
           If True, additional markup and tooltips are included in the produced
           PDF which indicate how tables and figures in the report correspond
           to members of this Results object.  These "tips" can be useful if
           you want to further manipulate the data contained in a table or
           figure.

        verbosity : int, optional
           How much detail to send to stdout.

        comm : mpi4py.MPI.Comm, optional
            When not None, an MPI communicator for distributing the computation
            across multiple processors.

        Returns
        -------
        None
        """

        printer = VerbosityPrinter.build_printer(verbosity, comm=comm)

        assert(self._bEssentialResultsSet)
        self.confidence_level = confidenceLevel
        self._comm = comm
        v = verbosity # shorthand

        if tips:
            def tooltiptex(directive):
                return '\\pdftooltip{{\\color{blue}\\texttt{%s}}\\quad}' \
                    % directive + '{Access this information in pyGSTi via' \
                    + '<ResultsObject>%s}' % directive
        else:
            def tooltiptex(directive):
                return "" #tooltips disabled


        #Get report output filename
        default_dir  = self.parameters['defaultDirectory']
        default_base = self.parameters['defaultBasename']

        if filename != "auto":
            report_dir = _os.path.dirname(filename)
            report_base = _os.path.splitext( _os.path.basename(filename) )[0] + suffix
        else:
            cwd = _os.getcwd()
            report_dir  = default_dir  if (default_dir  is not None) else cwd
            report_base = default_base if (default_base is not None) else "GSTBrief"
            report_base += suffix

        if datasetLabel == "auto":
            if default_base is not None:
                datasetLabel = _latex.latex_escaped( default_base )
            else:
                datasetLabel = "$\\mathcal{D}$"

        if title == "auto": title = "Brief GST report for %s" % datasetLabel

        ######  Generate Report ######
        #Steps:
        # 1) generate latex tables
        # 2) generate plots
        # 3) populate template latex file => report latex file
        # 4) compile report latex file into PDF
        # 5) remove auxiliary files generated during compilation


        if self._LsAndGermInfoSet:
            baseStr_dict = self._getBaseStrDict()
            Ls = self.parameters['max length list']
            st = 1 if Ls[0] == 0 else 0 #start index: skip LGST column in plots
            goodnessOfFitSection = True
        else:
            goodnessOfFitSection = False

        #Note: for now, we assume the best gateset corresponds to the last L-value
        best_gs = self.gatesets['final estimate']
        obj = self.parameters['objective']

        qtys = {} # dictionary to store all latex strings to be inserted into report template
        qtys['title'] = title
        qtys['datasetLabel'] = datasetLabel
        qtys['settoggles'] =  "\\toggle%s{confidences}\n" % \
            ("false" if confidenceLevel is None else "true")
        qtys['settoggles'] += "\\toggle%s{goodnessSection}\n" % \
            ("true" if goodnessOfFitSection else "false")
        qtys['confidenceLevel'] = "%g" % confidenceLevel \
            if confidenceLevel is not None else "NOT-SET"
        qtys['objective'] = "$\\log{\\mathcal{L}}$" \
            if obj == "logl" else "$\\chi^2$"
        qtys['gofObjective'] = "$2\\Delta\\log{\\mathcal{L}}$" \
            if obj == "logl" else "$\\chi^2$"

        if self.options.errgen_type == "logTiG":
            qtys['errorgenformula'] = "$\hat{G} = G_{\mathrm{target}}e^{\mathbb{L}}$"
        elif self.options.errgen_type == "logG-logT":
            qtys['errorgenformula'] = "$\hat{G} = e^{\mathbb{L} + \log G_{\mathrm{target}}}$"
        else:
            qtys['errorgenformula'] = "???"


        if confidenceLevel is not None:
            cri = self._get_confidence_region(confidenceLevel)
            qtys['confidenceIntervalScaleFctr'] = "%.3g" % cri.intervalScaling
            qtys['confidenceIntervalNumNonGaugeParams'] = "%d" % cri.nNonGaugeParams
        else:
            cri = None
            qtys['confidenceIntervalScaleFctr'] = "NOT-SET"
            qtys['confidenceIntervalNumNonGaugeParams'] = "NOT-SET"

        pdfInfo = [('Author','pyGSTi'), ('Title', title),
                   ('Keywords', 'GST'), ('pyGSTi Version',_version.__version__),
                   ('opt_long_tables', self.options.long_tables),
                   ('opt_table_class', self.options.table_class),
                   ('opt_template_path', self.options.template_path),
                   ('opt_latex_cmd', self.options.latex_cmd) ]
        # Note: use definite ordering of parameters
        for key in sorted(list(self.parameters.keys())):
            pdfInfo.append( (key, self.parameters[key]) )
        qtys['pdfinfo'] = _to_pdfinfo( pdfInfo )

        #Get figure directory for figure generation *and* as a
        # scratch space for tables.
        D = report_base + "_files" #figure directory relative to reportDir
        if not _os.path.isdir( _os.path.join(report_dir,D)):
            _os.mkdir( _os.path.join(report_dir,D))

        # 1) get latex tables
        printer.log("*** Generating tables ***")

        std_tables = ('bestGatesetSpamTable',
                      'bestGatesetSpamParametersTable',
                      'bestGatesetGatesTable',
                      'bestGatesetDecompTable','bestGatesetRotnAxisTable',
                      'bestGatesetVsTargetTable',
                      'bestGatesetErrorGenTable')
        gof_tables = ('progressTable',)

        tables_to_compute = std_tables
        tables_to_blank = []

        if goodnessOfFitSection:
            tables_to_compute += gof_tables
        else:
            tables_to_blank += gof_tables

        for key in tables_to_compute:
            qtys[key] = self.tables.get(key, verbosity=printer - 1).render(
                'latex',longtables=self.options.long_tables, scratchDir=D,
                precision=self.options.precision,
                polarprecision=self.options.polar_precision,
                sciprecision=self.options.sci_precision)
            qtys["tt_"+key] = tooltiptex(".tables['%s']" % key)

        for key in tables_to_blank:
            qtys[key] = _generation.get_blank_table().render(
                'latex',longtables=self.options.long_tables)
            qtys["tt_"+key] = ""


        # 2) generate plots
        printer.log("*** Generating plots ***")

        if _matplotlib.is_interactive():
            _matplotlib.pyplot.ioff()
            bWasInteractive = True
        else: bWasInteractive = False

        #if goodnessOfFitSection:
        #    strs = ( self.gatestring_lists['prep fiducials'],
        #             self.gatestring_lists['effect fiducials'] )
        #    D = report_base + "_files" #figure directory relative to reportDir
        #    if not _os.path.isdir( _os.path.join(report_dir,D)):
        #        _os.mkdir( _os.path.join(report_dir,D))
        #
        #    #Chi2 or logl plot
        #    nPlots = 1
        #    if self.parameters['objective'] == "chi2":
        #        plotFnName,plotFnLatex = "Chi2", "$\chi^2$"
        #    elif self.parameters['objective'] == "logl":
        #        plotFnName,plotFnLatex = "LogL", "$\\log(\\mathcal{L})$"
        #    else:
        #        raise ValueError("Invalid objective value: %s"
        #                         % self.parameters['objective'])
        #
        #    if verbosity > 0:
        #        print " -- %s plots (%d): " % (plotFnName, nPlots),; _sys.stdout.flush()
        #
        #    if verbosity > 0:
        #        print "1 ",; _sys.stdout.flush()
        #    figkey = 'bestEstimateColorBoxPlot'
        #    figFilenm = "best%sBoxes.pdf" % plotFnName
        #    fig = self.figures.get(figkey, verbosity=v)
        #    fig.save_to(_os.path.join(report_dir, D, figFilenm))
        #    maxX = fig.get_extra_info()['nUsedXs']; maxY = fig.get_extra_info()['nUsedYs']
        #    maxW,maxH = 6.5,8.5 #max width and height of graphic in latex document (in inches)
        #
        #    if verbosity > 0:
        #        print ""; _sys.stdout.flush()
        #
        #    qtys[figkey]  = "\\includegraphics[width=%.2fin,height=%.2fin,keepaspectratio]{%s/%s}" % (maxW,maxH,D,figFilenm)
        #    qtys['tt_'+ figkey]  = tooltiptex(".figures['%s']" % figkey)

        if bWasInteractive:
            _matplotlib.pyplot.ion()

        # 3) populate template latex file => report latex file
        printer.log("*** Merging into template file ***")

        mainTexFilename = _os.path.join(report_dir, report_base + ".tex")
        appendicesTexFilename = _os.path.join(report_dir, report_base + "_appendices.tex")
        pdfFilename = _os.path.join(report_dir, report_base + ".pdf")

        mainTemplate = "brief_report_main.tex"
        self._merge_template(qtys, mainTemplate, mainTexFilename)

        # 4) compile report latex file into PDF
        printer.log("Latex file(s) successfully generated.  Attempting to compile with pdflatex...")
        cwd = _os.getcwd()
        if len(report_dir) > 0:
            _os.chdir(report_dir)

        try:
            self._compile_latex_report(report_dir, report_base,
                                       self.options.latex_call, printer)
        except _subprocess.CalledProcessError as e:
            printer.error("pdflatex returned code %d " % e.returncode +
                          "Check %s.log to see details." % report_base)
        finally:
            _os.chdir(cwd)

        return


    def create_presentation_pdf(self, confidenceLevel=None, filename="auto",
                              title="auto", datasetLabel="auto", suffix="",
                              debugAidsAppendix=False,
                              pixelPlotAppendix=False, whackamoleAppendix=False,
                              m=0, M=10, verbosity=0, comm=None):
        """
        Create a GST presentation (i.e. slides) using the beamer latex package.

        The slides can contain most (but not all) of the tables and figures
        from the "full" report but contain only minimal descriptive text.  This
        output if useful for those familiar with the GST full report who need
        to present the results in a projector-friendly format.

        Parameters
        ----------
        confidenceLevel : float, optional
           If not None, then the confidence level (between 0 and 100) used in
           the computation of confidence regions/intervals. If None, no
           confidence regions or intervals are computed.

        filename : string, optional
           The output filename where the presentation file(s) will be saved.
           Specifying "auto" will use the default directory and base name
           (specified in set_additional_info) if given, otherwise the file
           "GSTSlides.pdf" will be output to the current directoy.

        title : string, optional
           The title of the presentation.  "auto" uses a default title which
           specifyies the label of the dataset as well.

        datasetLabel : string, optional
           A label given to the dataset.  If set to "auto", then the label
           will be the base name of the dataset filename without extension
           (if given) or "$\\mathcal{D}$" (if not).

        suffix : string, optional
           A suffix to add to the end of the presentation filename.  Useful when
           filename is "auto" and you generate different presentations using
           the same dataset.

        debugAidsAppendix : bool, optional
           Whether to include the "debugging aids" appendix.  This
           appendix contains comparisons of GST and Direct-GST and small-
           eigenvalue error rates among other quantities potentially
           useful for figuring out why the GST estimate did not fit
           the data as well as expected.

        pixelPlotAppendix : bool, optional
           Whether to include the "pixel plots" appendix, which shows
           the goodness of fit, in the form of color box plots, for the
           intermediate iterations of the GST algortihm.

        whackamoleAppendix : bool, optional
           Whether to include the "whack-a-mole" appendix, which contains
           colr box plots showing the effect of reducing ("whacking") one
           particular part of the overall goodness of fit box plot.

        m, M : float, optional
           Minimum and Maximum values of the color scale used in the
           presentation's color box plots.

        verbosity : int, optional
           How much detail to send to stdout.

        comm : mpi4py.MPI.Comm, optional
            When not None, an MPI communicator for distributing the computation
            across multiple processors.


        Returns
        -------
        None
        """

        printer = VerbosityPrinter.build_printer(verbosity, comm=comm)

        assert(self._bEssentialResultsSet)
        self.confidence_level = confidenceLevel
        self._comm = comm
        v = verbosity # shorthand

        #Currently, no tooltip option for presentations
        def tooltiptex(directive):
            return "" #tooltips disabled

        #Get report output filename
        default_dir  = self.parameters['defaultDirectory']
        default_base = self.parameters['defaultBasename']

        if filename != "auto":
            report_dir = _os.path.dirname(filename)
            report_base = _os.path.splitext( _os.path.basename(filename) )[0] + suffix
        else:
            cwd = _os.getcwd()
            report_dir  = default_dir  if (default_dir  is not None) else cwd
            report_base = default_base if (default_base is not None) else "GSTSlides"
            report_base += suffix

        if datasetLabel == "auto":
            if default_base is not None:
                datasetLabel = _latex.latex_escaped( default_base )
            else:
                datasetLabel = "$\\mathcal{D}$"

        if title == "auto": title = "GST on %s" % datasetLabel

        ######  Generate Presentation ######
        #Steps:
        # 1) generate latex tables
        # 2) generate plots
        # 3) populate template latex file => report latex file
        # 4) compile report latex file into PDF
        # 5) remove auxiliary files generated during compilation

        #Note: for now, we assume the best gateset corresponds to the last L-value
        best_gs = self.gatesets['final estimate']

        if not self._LsAndGermInfoSet: #cannot create appendices which depend on this structure
            debugAidsAppendix = False
            pixelPlotAppendix = False
            whackamoleAppendix = False

        qtys = {} # dictionary to store all latex strings to be inserted into report template
        qtys['title'] = title
        qtys['datasetLabel'] = datasetLabel
        qtys['settoggles'] =  "\\toggle%s{confidences}\n" % \
            ("false" if confidenceLevel is None else "true" )
        qtys['settoggles'] += "\\toggle%s{LsAndGermsSet}\n" % \
            ("true" if self._LsAndGermInfoSet else "false" )
        qtys['settoggles'] += "\\toggle%s{debuggingaidsappendix}\n" % \
            ("true" if debugAidsAppendix else "false" )
        qtys['settoggles'] += "\\toggle%s{pixelplotsappendix}\n" % \
            ("true" if pixelPlotAppendix else "false" )
        qtys['settoggles'] += "\\toggle%s{whackamoleappendix}\n" % \
            ("true" if whackamoleAppendix else "false" )
        qtys['confidenceLevel'] = "%g" % confidenceLevel \
            if confidenceLevel is not None else "NOT-SET"
        qtys['objective'] = "$\\log{\\mathcal{L}}$" \
            if self.parameters['objective'] == "logl" else "$\\chi^2$"
        qtys['gofObjective'] = "$2\\Delta\\log{\\mathcal{L}}$" \
            if self.parameters['objective'] == "logl" else "$\\chi^2$"

        if self.options.errgen_type == "logTiG":
            qtys['errorgenformula'] = "$\hat{G} = G_{\mathrm{target}}e^{\mathbb{L}}$"
        elif self.options.errgen_type == "logG-logT":
            qtys['errorgenformula'] = "$\hat{G} = e^{\mathbb{L} + \log G_{\mathrm{target}}}$"
        else:
            qtys['errorgenformula'] = "???"

        if confidenceLevel is not None:
            cri = self._get_confidence_region(confidenceLevel)
            qtys['confidenceIntervalScaleFctr'] = "%.3g" % cri.intervalScaling
            qtys['confidenceIntervalNumNonGaugeParams'] = "%d" % cri.nNonGaugeParams
        else:
            cri = None
            qtys['confidenceIntervalScaleFctr'] = "NOT-SET"
            qtys['confidenceIntervalNumNonGaugeParams'] = "NOT-SET"

        pdfInfo = [('Author','pyGSTi'), ('Title', title),
                   ('Keywords', 'GST'), ('pyGSTi Version',_version.__version__),
                   ('opt_long_tables', self.options.long_tables),
                   ('opt_table_class', self.options.table_class),
                   ('opt_template_path', self.options.template_path),
                   ('opt_latex_cmd', self.options.latex_cmd) ]

        # Note: use definite ordering of parameters
        for key in sorted(list(self.parameters.keys())):
            pdfInfo.append( (key, self.parameters[key]) )
        qtys['pdfinfo'] = _to_pdfinfo( pdfInfo )


        #Get figure directory for figure generation *and* as a
        # scratch space for tables.
        D = report_base + "_files" #figure directory relative to reportDir
        if not _os.path.isdir( _os.path.join(report_dir,D)):
            _os.mkdir( _os.path.join(report_dir,D))


        # 1) get latex tables
        printer.log("*** Generating tables ***")

        std_tables =('targetSpamTable','targetGatesTable',
                     'datasetOverviewTable','bestGatesetSpamTable',
                     'bestGatesetSpamParametersTable',
                     'bestGatesetGatesTable','bestGatesetChoiTable',
                     'bestGatesetDecompTable','bestGatesetRotnAxisTable',
                     'bestGatesetVsTargetTable','bestGatesetErrorGenTable')
        ls_and_germs_tables = ('fiducialListTable','prepStrListTable',
                               'effectStrListTable','germListTable',
                               'progressTable')

        tables_to_compute = std_tables
        tables_to_blank = []

        if self._LsAndGermInfoSet:
            tables_to_compute += ls_and_germs_tables
        else:
            tables_to_blank += ls_and_germs_tables

        for key in tables_to_compute:
            qtys[key] = self.tables.get(key, verbosity=printer - 1).render(
                'latex',longtables=self.options.long_tables, scratchDir=D,
                precision=self.options.precision,
                polarprecision=self.options.polar_precision,
                sciprecision=self.options.sci_precision)
            qtys["tt_"+key] = tooltiptex(".tables['%s']" % key)

        for key in tables_to_blank:
            qtys[key] = _generation.get_blank_table().render(
                'latex',longtables=self.options.long_tables)
            qtys["tt_"+key] = ""


        # 2) generate plots
        printer.log("*** Generating plots ***")

        if _matplotlib.is_interactive():
            _matplotlib.pyplot.ioff()
            bWasInteractive = True
        else: bWasInteractive = False

        maxW,maxH = 4.0,3.0 #max width and height of graphic in latex presentation (in inches)
        maxHc = 2.5 #max height allowed for a figure with a caption (in inches)

        def incgr(figFilenm,W=None,H=None): #includegraphics "macro"
            if W is None: W = maxW
            if H is None: H = maxH
            return "\\includegraphics[width=%.2fin,height=%.2fin" % (W,H) + \
                ",keepaspectratio]{%s/%s}" % (D,figFilenm)

        def set_fig_qtys(figkey, figFilenm, v, W=None,H=None):
            fig = self.figures.get(figkey, verbosity=v)
            fig.save_to(_os.path.join(report_dir, D, figFilenm))
            qtys[figkey] = incgr(figFilenm,W,H)
            qtys['tt_' + figkey] = tooltiptex(".figures['%s']" % figkey)
            return fig

        #Chi2 or logl plots
        if self._LsAndGermInfoSet:
            baseStr_dict = self._getBaseStrDict()

            Ls = self.parameters['max length list']
            st = 1 if Ls[0] == 0 else 0 #start index: skip LGST col in box plots
            nPlots = (len(Ls[st:])-1)+1 if pixelPlotAppendix else 1

            if self.parameters['objective'] == "chi2":
                plotFnName,plotFnLatex = "Chi2", "$\chi^2$"
            elif self.parameters['objective'] == "logl":
                plotFnName,plotFnLatex = "LogL", "$\\log(\\mathcal{L})$"
            else:
                raise ValueError("Invalid objective value: %s"
                                 % self.parameters['objective'])

            printer.log(" -- %s plots (%d): " % (plotFnName, nPlots), end='')
            with printer.progress_logging(1):
                printer.show_progress(0, 1, prefix='', end='')
            fig = set_fig_qtys("bestEstimateColorBoxPlot",
                               "best%sBoxes.pdf" % plotFnName, printer - 1)
            maxX = fig.get_extra_info()['nUsedXs']
            maxY = fig.get_extra_info()['nUsedYs']

        else:
            for figkey in ["bestEstimateColorBoxPlot"]:
                qtys[figkey] = qtys["tt_"+figkey] = ""


        pixplots = ""
        if pixelPlotAppendix:
            paramListLength = len(self.parameters['max length list'])-1

            with printer.progress_logging(1):
                for i in range(st, paramListLength):
                    printer.show_progress(i, paramListLength, prefix='', end='')
                    #printer.log("%d " % (i-st+2), end='')

                    fig = self.figures.get("estimateForLIndex%dColorBoxPlot" % i,
                                           verbosity=printer - 1)
                    fig.save_to( _os.path.join(report_dir, D,
                                               "L%d_%sBoxes.pdf" % (i,plotFnName)) )
                    lx = fig.get_extra_info()['nUsedXs']
                    ly = fig.get_extra_info()['nUsedYs']

                    #scale figure size according to number of rows and columns+1
                    # (+1 for labels ~ another col) relative to initial plot
                    W = float(lx+1)/float(maxX+1) * maxW
                    H = float(ly)  /float(maxY)   * maxH

                    pixplots += "\n"
                    pixplots += "\\begin{frame}\n"
                    pixplots += "\\frametitle{Iteration %d ($L=%d$): %s values}\n" \
                        % (i, self.parameters['max length list'][i], plotFnLatex)
                    pixplots += "\\begin{figure}\n"
                    pixplots += "\\begin{center}\n"
                    #pixplots += "\\adjustbox{max height=\\dimexpr\\textheight-5.5cm\\relax, max width=\\textwidth}{"
                    pixplots += "\\includegraphics[width=%.2fin,height=%.2fin,keepaspectratio]{%s/L%d_%sBoxes.pdf}\n" % (W,H,D,i,plotFnName)
                    #FUTURE: add caption and conditional tooltip string?
                    pixplots += "\\end{center}\n"
                    pixplots += "\\end{figure}\n"
                    pixplots += "\\end{frame}\n"

        #Set template quantity (empty string if appendix disabled)
        qtys['intermediate_pixel_plot_slides'] = pixplots

        printer.log("")

        if debugAidsAppendix:
            #Direct-GST and deviation
            printer.log(" -- Direct-X plots (2)", end='')
            with printer.progress_logging(1):
                printer.show_progress(0, 2, prefix='', end='')
                fig = set_fig_qtys("directLongSeqGSTColorBoxPlot",
                               "directLongSeqGST%sBoxes.pdf" % plotFnName, printer - 1, H=maxHc)
                printer.show_progress(1, 2, prefix='', end='')
                fig = set_fig_qtys("directLongSeqGSTDeviationColorBoxPlot",
                                   "directLongSeqGSTDeviationBoxes.pdf", printer - 1, H=maxHc)
            printer.log('')

            #Small eigenvalue error rate
            printer.log(" -- Error rate plots...")
            fig = set_fig_qtys("smallEigvalErrRateColorBoxPlot",
                               "smallEigvalErrRateBoxes.pdf", printer - 1, H=maxHc)

        else:
            for figkey in ["directLongSeqGSTColorBoxPlot",
                           "directLongSeqGSTDeviationColorBoxPlot",
                           "smallEigvalErrRateColorBoxPlot"]:
                qtys[figkey] = qtys["tt_"+figkey] = ""



        whackamoleplots = ""
        if whackamoleAppendix:
            #Whack-a-mole plots for highest L of each length-1 germ
            highestL = self.parameters['max length list'][-1]
            allGateStrings = self.gatestring_lists['all']
            hammerWeight = 10.0
            len1Germs = [ g for g in self.gatestring_lists['germs']
                          if len(g) == 1 ]

            printer.log(" -- Whack-a-mole plots (%d): " % (2*len(len1Germs)), end='')

            with printer.progress_logging(1):
                for i,germ in enumerate(len1Germs):
                    printer.show_progress(i, len(len1Germs), prefix='', end='')

                    fig = self.figures.get("whack%sMoleBoxes" % germ[0],verbosity=printer - 1)
                    fig.save_to(_os.path.join(report_dir, D,"whack%sMoleBoxes.pdf"
                                              % germ[0]))

                    whackamoleplots += "\n"
                    whackamoleplots += "\\begin{frame}\n"
                    whackamoleplots += "\\frametitle{Whack-a-%s-mole plot for $\mathrm{%s}^{%d}$}" % (plotFnLatex,germ[0],highestL)
                    whackamoleplots += "\\begin{figure}\n"
                    whackamoleplots += "\\begin{center}\n"
                    #whackamoleplots += "\\adjustbox{max height=\\dimexpr\\textheight-5.5cm\\relax, max width=\\textwidth}{"
                    whackamoleplots += "\\includegraphics[width=%.2fin,height=%.2fin,keepaspectratio]{%s/whack%sMoleBoxes.pdf}\n" % (maxW,maxH,D,germ[0])
                    #FUTURE: add caption and conditional tooltip?
                    whackamoleplots += "\\end{center}\n"
                    whackamoleplots += "\\end{figure}\n"
                    whackamoleplots += "\\end{frame}\n"

            with printer.progress_logging(1):
                for i,germ in enumerate(len1Germs):
                    # printer.log("%d " % (len(len1Germs)+i+1), end='')
                    printer.show_progress(i, len(len1Germs), prefix='', end='')

                    fig = self.figures.get("whack%sMoleBoxesSummed" % germ[0],
                                           verbosity=printer - 1)
                    fig.save_to(_os.path.join(report_dir, D,"whack%sMoleBoxesSummed.pdf" % germ[0]))

                    whackamoleplots += "\n"
                    whackamoleplots += "\\begin{frame}\n"
                    whackamoleplots += "\\frametitle{Summed whack-a-%s-mole plot for $\mathrm{%s}^{%d}$}" % (plotFnLatex,germ[0],highestL)
                    whackamoleplots += "\\begin{figure}\n"
                    whackamoleplots += "\\begin{center}\n"
                    #whackamoleplots += "\\adjustbox{max height=\\dimexpr\\textheight-5.5cm\\relax, max width=\\textwidth}{"
                    whackamoleplots += "\\includegraphics[width=%.2fin,height=%.2fin,keepaspectratio]{%s/whack%sMoleBoxesSummed.pdf}\n" % (maxW,maxH,D,germ[0])
                    #FUTURE: add caption and conditional tooltip?
                    whackamoleplots += "\\end{center}\n"
                    whackamoleplots += "\\end{figure}\n"
                    whackamoleplots += "\\end{frame}\n"

            printer.log('')

        #Set template quantity (empty string if appendix disabled)
        qtys['whackamole_plot_slides'] = whackamoleplots

        if bWasInteractive:
            _matplotlib.pyplot.ion()


        # 3) populate template latex file => report latex file
        printer.log("*** Merging into template file ***")

        mainTexFilename = _os.path.join(report_dir, report_base + ".tex")
        pdfFilename = _os.path.join(report_dir, report_base + ".pdf")

        mainTemplate = "slides_main.tex"
        self._merge_template(qtys, mainTemplate, mainTexFilename)


        # 4) compile report latex file into PDF
        printer.log("Latex file(s) successfully generated.  Attempting to compile with pdflatex...")
        cwd = _os.getcwd()
        if len(report_dir) > 0:
            _os.chdir(report_dir)

        try:
            self._compile_latex_report(report_dir, report_base,
                                       self.options.latex_call, printer)
        except _subprocess.CalledProcessError as e:
            printer.error("pdflatex returned code %d " % e.returncode +
                          "Check %s.log to see details." % report_base)
        finally:
            _os.chdir(cwd)

        return




    def create_presentation_ppt(self, confidenceLevel=None, filename="auto",
                            title="auto", datasetLabel="auto", suffix="",
                            debugAidsAppendix=False,
                            pixelPlotAppendix=False, whackamoleAppendix=False,
                            m=0, M=10, verbosity=0, pptTables=False, comm=None):
        """
        Create a GST Microsoft Powerpoint presentation.

        These slides can contain most (but not all) of the tables and figures
        from the "full" report but contain only minimal descriptive text.  This
        method uses the python-pptx package to write Powerpoint files.  The
        resulting powerpoint slides are meant to parallel those of the PDF
        presentation but are not as nice and clean.  This method exists because
        the Powerpoint format is an industry standard and makes it very easy to
        shamelessly co-opt GST figures or entire slides for incorporation into
        other presentations.

        Parameters
        ----------
        confidenceLevel : float, optional
           If not None, then the confidence level (between 0 and 100) used in
           the computation of confidence regions/intervals. If None, no
           confidence regions or intervals are computed.

        filename : string, optional
           The output filename where the presentation file(s) will be saved.
           Specifying "auto" will use the default directory and base name
           (specified in set_additional_info) if given, otherwise the file
           "GSTSlides.pptx" will be output to the current directoy.

        title : string, optional
           The title of the presentation.  "auto" uses a default title which
           specifyies the label of the dataset as well.

        datasetLabel : string, optional
           A label given to the dataset.  If set to "auto", then the label
           will be the base name of the dataset filename without extension
           (if given) or "D" (if not).

        suffix : string, optional
           A suffix to add to the end of the presentation filename.  Useful when
           filename is "auto" and you generate different presentations using
           the same dataset.

        debugAidsAppendix : bool, optional
           Whether to include the "debugging aids" appendix.  This
           appendix contains comparisons of GST and Direct-GST and small-
           eigenvalue error rates among other quantities potentially
           useful for figuring out why the GST estimate did not fit
           the data as well as expected.

        pixelPlotAppendix : bool, optional
           Whether to include the "pixel plots" appendix, which shows
           the goodness of fit, in the form of color box plots, for the
           intermediate iterations of the GST algortihm.

        whackamoleAppendix : bool, optional
           Whether to include the "whack-a-mole" appendix, which contains
           colr box plots showing the effect of reducing ("whacking") one
           particular part of the overall goodness of fit box plot.

        m, M : float, optional
           Minimum and Maximum values of the color scale used in the
           presentation's color box plots.

        verbosity : int, optional
           How much detail to send to stdout.

        pptTables : bool, optional
           If True, native powerpoint-format tables are placed in slides instead
           of creating the prettier-looking PNG images of latex-ed tables (which
           are used when False).  This option can be useful when you want to
           modify or extract a part of a table.

        comm : mpi4py.MPI.Comm, optional
            When not None, an MPI communicator for distributing the computation
            across multiple processors.


        Returns
        -------
        None
        """

        printer = VerbosityPrinter.build_printer(verbosity, comm=comm)

        assert(self._bEssentialResultsSet)
        self.confidence_level = confidenceLevel
        self._comm = comm
        v = verbosity # shorthand

        #Currently, no tooltip option for presentations
        def tooltiptext(directive):
            return "" #tooltips disabled


        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
        except:
            raise ValueError("Cannot import pptx: it seems like python-pptx is not installed on your system!")

        try:
            from PIL import Image
        except:
            raise ValueError("Cannot import PIL: it seems like the python imaging library is not installed on your system!")

        #Get report output filename
        default_dir  = self.parameters['defaultDirectory']
        default_base = self.parameters['defaultBasename']

        if filename != "auto":
            report_dir = _os.path.dirname(filename)
            report_base = _os.path.splitext( _os.path.basename(filename) )[0] + suffix
        else:
            cwd = _os.getcwd()
            report_dir  = default_dir  if (default_dir  is not None) else cwd
            report_base = default_base if (default_base is not None) else "GSTSlides"
            report_base += suffix

        if datasetLabel == "auto":
            if default_base is not None:
                datasetLabel =  default_base
            else:
                datasetLabel = "D"

        if title == "auto": title = "GST on %s" % datasetLabel

        ######  Generate PPT Presentation ######

        #Note: for now, we assume the best gateset corresponds to the last L-value
        best_gs = self.gatesets['final estimate']

        if not self._LsAndGermInfoSet: #cannot create appendices which depend on this structure
            debugAidsAppendix = False
            pixelPlotAppendix = False
            whackamoleAppendix = False

        qtys = {}
        qtys['title'] = title
        qtys['datasetLabel'] = datasetLabel
        qtys['confidenceLevel'] = "%g" % confidenceLevel \
            if confidenceLevel is not None else "NOT-SET"
        qtys['objective'] = "$\\log{\\mathcal{L}}$" \
            if self.parameters['objective'] == "logl" else "$\\chi^2$"
        qtys['gofObjective'] = "$2\\Delta\\log{\\mathcal{L}}$" \
            if self.parameters['objective'] == "logl" else "$\\chi^2$"

        if confidenceLevel is not None:
            cri = self._get_confidence_region(confidenceLevel)
            qtys['confidenceIntervalScaleFctr'] = "%.3g" % cri.intervalScaling
            qtys['confidenceIntervalNumNonGaugeParams'] = "%d" % cri.nNonGaugeParams
        else:
            cri = None
            qtys['confidenceIntervalScaleFctr'] = "NOT-SET"
            qtys['confidenceIntervalNumNonGaugeParams'] = "NOT-SET"

        #Get figure directory for figure generation *and* as a
        # scratch space for tables.
        D = report_base + "_files" #figure directory relative to reportDir
        if not _os.path.isdir( _os.path.join(report_dir,D)):
            _os.mkdir( _os.path.join(report_dir,D))

        # 1) get ppt tables
        printer.log("*** Generating tables ***")

        std_tables = ('targetSpamTable','targetGatesTable',
                      'datasetOverviewTable', 'bestGatesetSpamTable',
                      'bestGatesetSpamParametersTable','bestGatesetGatesTable',
                      'bestGatesetChoiTable', 'bestGatesetDecompTable',
                      'bestGatesetRotnAxisTable','bestGatesetVsTargetTable',
                      'bestGatesetErrorGenTable')

        ls_and_germs_tables = ('fiducialListTable','prepStrListTable',
                               'effectStrListTable','germListTable',
                               'progressTable')

        tables_to_compute = std_tables
        tables_to_blank = []

        if self._LsAndGermInfoSet:
            tables_to_compute += ls_and_germs_tables
        else:
            tables_to_blank += ls_and_germs_tables

        with printer.progress_logging(1):
            for i, key in enumerate(tables_to_compute):
                printer.show_progress(i, len(tables_to_compute), prefix='', end='')
                qtys[key] = self.tables.get(key, verbosity=printer)
                qtys["tt_"+key] = tooltiptext(".tables['%s']" % key)

        for key in tables_to_blank:
            qtys[key] = _generation.get_blank_table()
            qtys["tt_"+key] = ""


        # 2) generate plots
        printer.log("*** Generating plots ***")

        if _matplotlib.is_interactive():
            _matplotlib.pyplot.ioff()
            bWasInteractive = True
        else: bWasInteractive = False

        fileDir = _os.path.join(report_dir, D)
        maxW,maxH = 4.0,3.0 #max width and height of graphic in latex presentation (in inches)

        def incgr(figFilenm,W=None,H=None): #includegraphics "macro"
            return "%s/%s" % (fileDir,figFilenm)

        def set_fig_qtys(figkey, figFilenm, v, W=None,H=None):
            fig = self.figures.get(figkey, verbosity=v)
            fig.save_to(_os.path.join(report_dir, D, figFilenm))
            qtys[figkey] = incgr(figFilenm,W,H)
            qtys['tt_' + figkey] = tooltiptext(".figures['%s']" % figkey)
            return fig


        #Chi2 or logl plots
        if self._LsAndGermInfoSet:
            baseStr_dict = self._getBaseStrDict()

            Ls = self.parameters['max length list']
            st = 1 if Ls[0] == 0 else 0 #start index: skip LGST col in box plots
            nPlots = (len(Ls[st:])-1)+1 if pixelPlotAppendix else 1

            if self.parameters['objective'] == "chi2":
                plotFnName,plotFnLatex = "Chi2", "$\chi^2$"
            elif self.parameters['objective'] == "logl":
                plotFnName,plotFnLatex = "LogL", "$\\log(\\mathcal{L})$"
            else:
                raise ValueError("Invalid objective value: %s" \
                                     % self.parameters['objective'])

            printer.log(" -- %s plots (%d): " % (plotFnName, nPlots), end='')
            with printer.progress_logging(1):
                printer.show_progress(0, 1, prefix='', end='')
            fig = set_fig_qtys("bestEstimateColorBoxPlot",
                               "best%sBoxes.png" % plotFnName, printer - 1)
            maxX = fig.get_extra_info()['nUsedXs']
            maxY = fig.get_extra_info()['nUsedYs']

        else:
            for figkey in ["bestEstimateColorBoxPlot"]:
                qtys[figkey] = qtys["tt_"+figkey] = ""


        pixplots = []
        if pixelPlotAppendix:
            Ls = self.parameters['max length list']
            with printer.progress_logging(1):
                for i in range(st,len(Ls)-1):

                    printer.show_progress(i, len(Ls)-1, prefix='', end='')
                    # printer.log("%d " % (i-st+2), end='')

                    fig = self.figures.get("estimateForLIndex%dColorBoxPlot" % i,
                                           verbosity=printer - 1)
                    fig.save_to( _os.path.join(report_dir, D,"L%d_%sBoxes.png" %
                                               (i,plotFnName)) )
                    lx = fig.get_extra_info()['nUsedXs']
                    ly = fig.get_extra_info()['nUsedYs']

                    #scale figure size according to number of rows and columns+1
                    # (+1 for labels ~ another col) relative to initial plot
                    W = float(lx+1)/float(maxX+1) * maxW
                    H = float(ly)  /float(maxY)   * maxH

                    pixplots.append( _os.path.join(
                            report_dir, D, "L%d_%sBoxes.png" % (i,plotFnName)) )
                    #FUTURE: Add tooltip caption info further down?

        #Set template quantity (empty array if appendix disabled)
        qtys['intermediate_pixel_plot_slides'] = pixplots

        printer.log("")

        if debugAidsAppendix:
            #Direct-GST and deviation
            printer.log(" -- Direct-X plots (2)", end="")
            with printer.progress_logging(1):
                printer.show_progress(0, 2, prefix='', end='')
                fig = set_fig_qtys("directLongSeqGSTColorBoxPlot",
                               "directLongSeqGST%sBoxes.png" % plotFnName, printer - 1)

                printer.show_progress(1, 2, prefix='', end='')
                fig = set_fig_qtys("directLongSeqGSTDeviationColorBoxPlot",
                                   "directLongSeqGSTDeviationBoxes.png", printer - 1)

                printer.log('')
                #Small eigenvalue error rate
                printer.log(" -- Error rate plots...")
                fig = set_fig_qtys("smallEigvalErrRateColorBoxPlot",
                                   "smallEigvalErrRateBoxes.png", printer - 1)

        else:
            for figkey in ["directLongSeqGSTColorBoxPlot",
                           "directLongSeqGSTDeviationColorBoxPlot",
                           "smallEigvalErrRateColorBoxPlot"]:
                qtys[figkey] = qtys["tt_"+figkey] = ""


        whackamoleplots = []
        if whackamoleAppendix:
            #Whack-a-mole plots for highest L of each length-1 germ
            Ls = self.parameters['max length list']
            highestL = Ls[-1]; allGateStrings = self.gatestring_lists['all']
            hammerWeight = 10.0
            len1Germs = [ g for g in self.gatestring_lists['germs']
                          if len(g) == 1 ]

            printer.log(" -- Whack-a-mole plots (%d): " % (2*len(len1Germs)), end='')
            with printer.progress_logging(1):
                for i,germ in enumerate(len1Germs):
                    printer.show_progress(i, len(len1Germs), prefix='', end='')

                    fig = self.figures.get("whack%sMoleBoxes" % germ[0],verbosity=printer -1)
                    fig.save_to(_os.path.join(report_dir, D,"whack%sMoleBoxes.png"
                                              % germ[0]))
                    whackamoleplots.append( _os.path.join(
                            report_dir, D, "whack%sMoleBoxes.png" % germ[0]) )
                    #FUTURE: Add tooltip caption info further down?

            with printer.progress_logging(1):
                for i,germ in enumerate(len1Germs):
                    printer.show_progress(i, len(len1Germs), prefix='', end='')
                    # printer.log("%d " % (len(len1Germs)+i+1), end='')

                    fig = self.figures.get("whack%sMoleBoxesSummed" % germ[0],
                                           verbosity=printer - 1)
                    fig.save_to(_os.path.join(
                            report_dir, D, "whack%sMoleBoxesSummed.png" % germ[0]))
                    whackamoleplots.append( _os.path.join(
                            report_dir, D,"whack%sMoleBoxesSummed.png" % germ[0]) )
                    #FUTURE: Add tooltip caption info further down?

            printer.log('')

        #Set template quantity (empty array if appendix disabled)
        qtys['whackamole_plot_slides'] = whackamoleplots

        if bWasInteractive:
            _matplotlib.pyplot.ion()


        # 3) create PPT file via python-pptx
        printer.log("*** Assembling PPT file ***")

        mainPPTFilename = _os.path.join(report_dir, report_base + ".pptx")
        templatePath = self.options.template_path \
            if (self.options.template_path is not None) else \
            _os.path.join( _os.path.dirname(_os.path.abspath(__file__)),
                           "templates")
        templateFilename = _os.path.join( templatePath, "slides_main.pptx" )

        # slide layout indices

          #For normal powerpoint templates
        #SLD_LAYOUT_TITLE = 0
        #SLD_LAYOUT_TITLE_AND_CONTENT = 1
        #SLD_LAYOUT_TITLE_NO_CONTENT = 5

          #For sandia's template
        SLD_LAYOUT_TITLE = 0
        SLD_LAYOUT_TITLE_AND_CONTENT = 6
        SLD_LAYOUT_TITLE_NO_CONTENT = 10

        def draw_table_ppt(shapes, key, left, top, width, height, ptSize=10):
            tabDict = qtys[key].render('ppt')
            nRows = len(tabDict['row data']) + 1
            nCols = len(tabDict['column names'])

            def maxllen(s): #maximum line length
                return max( [len(ln) for ln in s.split('\n')] )

            max_col_widths = []
            for i in range(nCols):
                max_col_width = max( [ maxllen(rd[i]) for rd in tabDict['row data'] ] )
                max_col_width = max(max_col_width, maxllen(tabDict['column names'][i]) )
                max_col_widths.append(max_col_width)

            max_row_heights = []
            max_row_heights.append( max([ len(nm.split('\n')) for nm in tabDict['column names'] ] ) )
            for row_data in tabDict['row data']:
                max_row_heights.append( max([ len(el.split('\n')) for el in row_data ] ) )

            l = Inches(left); t = Inches(top)
            maxW = Inches(width); maxH = Inches(height)

            #table width and height per "pt" - scaling factors determined by trial and error
            Wfctr = 1.2; Hfctr = 1.2
            WperPt = Wfctr * sum(max_col_widths)
            HperPt = Hfctr * sum(max_row_heights)

            #set *desired* font size
            fontSize = Pt(ptSize)

            #compute width and height of table based on desired values
            w = int(WperPt * fontSize)
            h = int(HperPt * fontSize)

            #if table is too big, reduce point size until it fits
            if w > maxW:
                fontSize = maxW / WperPt
                w = maxW
                h = int(HperPt * fontSize)

            if h > maxH:
                fontSize = maxH / HperPt
                h = maxH
                w = int(WperPt * fontSize)

            table = shapes.add_table(nRows, nCols, l, t, w, h).table
            table.first_row = True
            table.horz_banding = True

            # set column widths
            for i in range(nCols):
                table.columns[i].width = int(max_col_widths[i] * Wfctr * fontSize)

            # set row heights
            for i in range(nRows):
                table.rows[i].height = int(max_row_heights[i] * Hfctr * fontSize)

            # write column headings
            for i in range(nCols):
                table.cell(0, i).text = tabDict['column names'][i]
                tf = table.cell(0, i).text_frame
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                tf.word_wrap = False
                tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                for run in table.cell(0, i).text_frame.paragraphs[0].runs:
                    run.font.bold = True
                    run.font.size = fontSize

            # write body cells
            for i in range(1,nRows):
                for j in range(nCols):
                    table.cell(i, j).text = tabDict['row data'][i-1][j]
                    tf = table.cell(i, j).text_frame
                    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                    tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
                    tf.word_wrap = False
                    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                    for run in table.cell(i, j).text_frame.paragraphs[0].runs:
                        run.font.size = fontSize

            return table

        def draw_table_latex(shapes, key, left, top, width, height, ptSize=10):
            latexTabStr = qtys[key].render(
                'latex', longtables=self.options.long_tables, scratchDir=D,
                precision=self.options.precision,
                polarprecision=self.options.polar_precision,
                sciprecision=self.options.sci_precision)
            d = {'toLatex': latexTabStr }
            printer.log("Latexing %s table..." % key)
            outputFilename = _os.path.join(fileDir, "%s.tex" % key)
            self._merge_template(d, "standalone.tex", outputFilename)

            cwd = _os.getcwd()
            _os.chdir(fileDir)

            try:
                latex_cmd = self.options.latex_call + \
                            ["-shell-escape", "%s.tex" % key]
                stdout, stderr, returncode = self._process_call(latex_cmd)
                self._evaluate_call(latex_cmd, stdout, stderr, returncode,
                                    printer)
                # Check to see if the PNG was generated
                if not _os.path.isfile("%s.png" % key):
                    raise Exception("File %s.png was not created by pdflatex"
                                    % key)
                _os.remove( "%s.tex" % key )
                _os.remove( "%s.log" % key )
                _os.remove( "%s.aux" % key )
            except _subprocess.CalledProcessError as e:
                printer.error("pdflatex returned code %d " % e.returncode +
                              "trying to render standalone %s. " % key +
                              "Check %s.log to see details." % key)
            except:
                printer.error("pdflatex failed to render standalone %s" % key)
                raise

            finally:
                _os.chdir(cwd)

            pathToImg = _os.path.join(fileDir, "%s.png" % key)
            return draw_pic(shapes, pathToImg, left, top, width, height)


        def draw_pic(shapes, path, left, top, width, height):
            with open(path, 'rb') as imagefile:
                pxWidth, pxHeight = Image.open(imagefile).size
            pxAspect = pxWidth / float(pxHeight) #aspect ratio of image
            maxAspect = width / float(height) #aspect ratio of "max" box
            if pxAspect > maxAspect:
                w = Inches(width); h = None # image is wider & flatter than max box => constrain width
                #print "%s -> constrain width to %f so height is %f" % (path,width,width / pxAspect)
            else:
                h = Inches(height); w = None # image is taller & thinner than max box => constrain height
                #print "%s -> constrain height to %f so width is %f" % (path,height,height * pxAspect)
            return shapes.add_picture(path, Inches(left), Inches(top), w, h)

        def add_slide(typ, title):
            slide_layout = prs.slide_layouts[typ]
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            return slide

        def add_text(shapes, left, top, width, height, text, ptSize=20):
            return add_text_list(shapes, left, top, width, height, [text])

        def add_text_list(shapes, left, top, width, height, textList, ptSize=20):
            l = Inches(left); t = Inches(top); w = Inches(width); h = Inches(height)
            txtBox = slide.shapes.add_textbox(l, t, w, h)
            tf = txtBox.text_frame
            tf.text = textList[0]
            for run in tf.paragraphs[0].runs:
                run.font.size = Pt(ptSize)

            for line in textList[1:]:
                p = tf.add_paragraph()
                p.text = line
                for run in p.runs:
                    run.font.size = Pt(ptSize)

            return txtBox

        drawTable = draw_table_ppt if pptTables else draw_table_latex

        # begin presentation creation
        #prs = Presentation() #templateFilename)
        templateDir =_os.path.join(_os.path.dirname(_os.path.abspath(__file__)),
                                  "templates") if self.options.template_path \
                                  is None else self.options.template_path
        prs = Presentation( _os.path.join( templateDir, "GSTTemplate.pptx" ) )


        # title slide
        slide = add_slide(SLD_LAYOUT_TITLE, title)
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = "Your GST results in Powerpoint!"

        # goodness of fit
        if self._LsAndGermInfoSet:
            slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "%s vs GST iteration" % plotFnName)
            #body_shape = slide.shapes.placeholders[1]; tf = body_shape.text_frame
            add_text_list(slide.shapes, 1, 2, 8, 2, ['Ns is the number of gate strings', 'Np is the number of parameters'], 15)
            drawTable(slide.shapes, 'progressTable', 1, 3, 8.5, 4, ptSize=10)

            slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Detailed %s Analysis" % plotFnName)
            draw_pic(slide.shapes, qtys['bestEstimateColorBoxPlot'], 1, 1.5, 8, 5.5)

        # gate esimtates
        slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "GST Estimate vs. target")
        drawTable(slide.shapes, 'bestGatesetVsTargetTable', 1.5, 1.5, 7, 2, ptSize=10)
        drawTable(slide.shapes, 'bestGatesetErrorGenTable', 1.5, 3.7, 7, 3.5, ptSize=10)

        slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "GST Estimate decomposition")
        drawTable(slide.shapes, 'bestGatesetDecompTable', 1, 1.5, 7.5, 3.5 , ptSize=10)
        drawTable(slide.shapes, 'bestGatesetRotnAxisTable', 1, 5.1, 5, 1.5, ptSize=10)

        slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Raw GST Estimate: Gates")
        drawTable(slide.shapes, 'bestGatesetGatesTable', 1, 2, 8, 5, ptSize=10)

        slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Raw GST Estimate: SPAM")
        drawTable(slide.shapes, 'bestGatesetSpamTable', 1, 1.5, 8, 3, ptSize=10)
        drawTable(slide.shapes, 'bestGatesetSpamParametersTable', 1, 5, 3, 1, ptSize=10)

        slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Raw GST Choi Matrices")
        drawTable(slide.shapes, 'bestGatesetChoiTable', 0.5, 1.5, 8.5, 5.5, ptSize=10)

        #Inputs to GST
        slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Target SPAM")
        drawTable(slide.shapes, 'targetSpamTable', 1.5, 1.5, 7, 5, ptSize=10)

        slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Target Gates")
        drawTable(slide.shapes, 'targetGatesTable', 1, 1.5, 7, 5, ptSize=10)

        if self._LsAndGermInfoSet:
            slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Fiducial and Germ Gate Strings")
            drawTable(slide.shapes, 'fiducialListTable', 1, 1.5, 4, 3, ptSize=10)
            drawTable(slide.shapes, 'germListTable', 5.5, 1.5, 4, 5, ptSize=10)

        slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Dataset Overview")
        drawTable(slide.shapes, 'datasetOverviewTable', 1, 2, 5, 4, ptSize=10)

        if debugAidsAppendix:
            #Debugging aids slides
            slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Direct-GST")
            draw_pic(slide.shapes, qtys['directLongSeqGSTColorBoxPlot'], 1, 1.5, 8, 5.5)

            slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Direct-GST Deviation")
            draw_pic(slide.shapes, qtys['directLongSeqGSTDeviationColorBoxPlot'], 1, 1.5, 8, 5.5)

            slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Per-gate error rates")
            draw_pic(slide.shapes, qtys['smallEigvalErrRateColorBoxPlot'], 1, 1.5, 8, 5.5)

        if pixelPlotAppendix:
            Ls = self.parameters['max length list']
            for i,pixPlotPath in zip( list(range(st,len(Ls)-1)), pixplots ):
                slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Iteration %d (L=%d): %s values" % (i,Ls[i],plotFnName))
                draw_pic(slide.shapes, pixPlotPath, 1, 1.5, 8, 5.5)

        if whackamoleAppendix:
            curPlot = 0
            for i,germ in enumerate(len1Germs):
                slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Whack-a-%s-mole plot for %s^%d" % (plotFnName,germ[0],highestL))
                draw_pic(slide.shapes, whackamoleplots[curPlot], 1, 1.5, 8, 5.5)
                curPlot += 1

            for i,germ in enumerate(len1Germs):
                slide = add_slide(SLD_LAYOUT_TITLE_NO_CONTENT, "Summed whack-a-%s-mole plot for %s^%d" % (plotFnName,germ[0],highestL))
                draw_pic(slide.shapes, whackamoleplots[curPlot], 1, 1.5, 8, 5.5)
                curPlot += 1

        # 4) save presenation as PPTX file
        prs.save(mainPPTFilename)
        printer.log("Final output PPT %s successfully generated." % mainPPTFilename)
        return


    def create_general_report_pdf(self, confidenceLevel=None, filename="auto",
                                  title="auto", datasetLabel="auto", suffix="",
                                  tips=False, verbosity=0, comm=None,
                                  showAppendix=False):
        """
        Create a "general" GST report.  This report is suited to display
        results for any number of qubits, and is detailed in the sense that
        it includes background and explanation text to help the user
        interpret the contained results.

        Parameters
        ----------
        confidenceLevel : float, optional
           If not None, then the confidence level (between 0 and 100) used in
           the computation of confidence regions/intervals. If None, no
           confidence regions or intervals are computed.

        filename : string, optional
           The output filename where the report file(s) will be saved.  Specifying
           "auto" will use the default directory and base name (specified in
           set_additional_info) if given, otherwise the file "GSTReport.pdf" will
           be output to the current directoy.  If None, then results are computed
           but no file output is generated (useful for pre-computing cached
           derived quantities).

        title : string, optional
           The title of the report.  "auto" uses a default title which
           specifyies the label of the dataset as well.

        datasetLabel : string, optional
           A label given to the dataset.  If set to "auto", then the label
           will be the base name of the dataset filename without extension
           (if given) or "$\\mathcal{D}$" (if not).

        suffix : string, optional
           A suffix to add to the end of the report filename.  Useful when
           filename is "auto" and you generate different reports using
           the same dataset.

        tips : boolean, optional
           If True, additional markup and tooltips are included in the produced
           PDF which indicate how tables and figures in the report correspond
           to members of this Results object.  These "tips" can be useful if
           you want to further manipulate the data contained in a table or
           figure.

        verbosity : int, optional
           How much detail to send to stdout.

        comm : mpi4py.MPI.Comm, optional
            When not None, an MPI communicator for distributing the computation
            across multiple processors.

        showAppendix : bool, optional
            Whether to display the appendix.

        Returns
        -------
        None
        """

        printer = VerbosityPrinter.build_printer(verbosity, comm=comm)
        tStart = _time.time()

        assert(self._bEssentialResultsSet)
        self.confidence_level = confidenceLevel
        self._comm = comm
          #set "current" level, used by ResultCache member dictionaries

        if self.parameters['objective'] != "logl":
            raise NotImplementedError("General reports are currently " +
                                      "only implemented for log-likelihood " +
                                      "objective function case.")

        if tips:
            def tooltiptex(directive):
                return '\\pdftooltip{{\\color{blue}\\texttt{%s}}\\quad}' \
                    % directive + '{Access this information in pyGSTi via' \
                    + '<ResultsObject>%s}' % directive


        else:
            def tooltiptex(directive):
                return "" #tooltips disabled


        #Get report output filename
        default_dir = self.parameters['defaultDirectory']
        default_base = self.parameters['defaultBasename']
        bOutputFiles = bool(filename is not None)

        if filename == "auto":
            cwd = _os.getcwd()
            report_dir  = default_dir  if (default_dir  is not None) else cwd
            report_base = default_base if (default_base is not None) \
                                       else "GSTReport"
            report_base += suffix
        elif filename is not None:
            report_dir = _os.path.dirname(filename)
            report_base = _os.path.splitext( _os.path.basename(filename) )[0] \
                           + suffix
        else:
            report_dir = None
            report_base = None

        if datasetLabel == "auto":
            if default_base is not None:
                datasetLabel = _latex.latex_escaped( default_base )
            else:
                datasetLabel = "$\\mathcal{D}$"

        if title == "auto": title = "GST report for %s" % datasetLabel

        ######  Generate Report ######
        #Steps:
        # 1) generate latex tables
        # 2) generate plots
        # 3) populate template latex file => report latex file
        # 4) compile report latex file into PDF
        # 5) remove auxiliary files generated during compilation
        #  FUTURE?? determine what we need to compute & plot by reading
        #           through the template file?

        #Note: for now, we assume the best gateset corresponds to the last
        #      L-value
        best_gs = self.gatesets['final estimate']
        v = verbosity # shorthand

        if not self._LsAndGermInfoSet: #cannot create appendices
            debugAidsAppendix = False  # which depend on this structure
            pixelPlotAppendix = False
            whackamoleAppendix = False

        qtys = {} # dictionary to store all latex strings
                  # to be inserted into report template
        qtys['title'] = title
        qtys['datasetLabel'] = datasetLabel
        qtys['settoggles'] =  "\\toggle%s{confidences}\n" % \
            ("false" if confidenceLevel is None else "true")
        qtys['settoggles'] += "\\toggle%s{LsAndGermsSet}\n" % \
            ("true" if self._LsAndGermInfoSet else "false")
        qtys['settoggles'] += "\\toggle%s{showAppendix}\n" % \
            ("true" if showAppendix else "false")
        #qtys['settoggles'] += "\\toggle%s{gaugeoptappendix}\n" % \
        #    ("true" if gaugeOptAppendix else "false")
        #qtys['settoggles'] += "\\toggle%s{pixelplotsappendix}\n" % \
        #    ("true" if pixelPlotAppendix else "false")
        #qtys['settoggles'] += "\\toggle%s{whackamoleappendix}\n" % \
        #    ("true" if whackamoleAppendix else "false")
        qtys['confidenceLevel'] = "%g" % \
            confidenceLevel if confidenceLevel is not None else "NOT-SET"
        qtys['linlg_pcntle'] = self.parameters['linlogPercentile']

        if confidenceLevel is not None:
            cri = self._get_confidence_region(confidenceLevel)
            qtys['confidenceIntervalScaleFctr'] = \
                "%.3g" % cri.intervalScaling
            qtys['confidenceIntervalNumNonGaugeParams'] = \
                "%d" % cri.nNonGaugeParams
        else:
            cri = None
            qtys['confidenceIntervalScaleFctr'] = "NOT-SET"
            qtys['confidenceIntervalNumNonGaugeParams'] = "NOT-SET"

        if self.options.errgen_type == "logTiG":
            qtys['errorgenformula'] = "$\hat{G} = G_{\mathrm{target}}e^{\mathbb{L}}$"
        elif self.options.errgen_type == "logG-logT":
            qtys['errorgenformula'] = "$\hat{G} = e^{\mathbb{L} + \log G_{\mathrm{target}}}$"
        else:
            qtys['errorgenformula'] = "???"


        pdfInfo = [('Author','pyGSTi'), ('Title', title),
                   ('Keywords', 'GST'), ('pyGSTi Version',_version.__version__),
                   ('opt_long_tables', self.options.long_tables),
                   ('opt_table_class', self.options.table_class),
                   ('opt_template_path', self.options.template_path),
                   ('opt_latex_cmd', self.options.latex_cmd) ]

        # Note: use definite ordering of parameters
        for key in sorted(list(self.parameters.keys())):
            pdfInfo.append( (key, self.parameters[key]) )
        qtys['pdfinfo'] = _to_pdfinfo( pdfInfo )

        #Get figure directory for figure generation *and* as a
        # scratch space for tables.
        if bOutputFiles:
            D = report_base + "_files" #figure directory relative to reportDir
            if not _os.path.isdir( _os.path.join(report_dir,D)):
                _os.mkdir( _os.path.join(report_dir,D))
        else:
            D = "" #empty scratch dir for figures

        # 1) get latex tables
        printer.log("*** Generating tables *** (%.1fs elapsed)"
                    % (_time.time()-tStart))


        std_tables = \
            ('targetSpamBriefTable', 'bestGatesetSpamBriefTable',
             'bestGatesetSpamParametersTable', 'bestGatesetVsTargetTable',
             'bestGatesetSpamVsTargetTable', 'bestGatesetGaugeOptParamsTable',
             'bestGatesetChoiEvalTable', 'datasetOverviewTable',
             'bestGatesetEvalTable', 'bestGatesetRelEvalTable',
             'targetGatesBoxTable', 'bestGatesetGatesBoxTable',
             'bestGatesetErrGenBoxTable')

        ls_and_germs_tables = ('fiducialListTable','prepStrListTable',
                               'effectStrListTable','germList2ColTable',
                               'progressTable')

        appendix_tables = ('bestGatesetErrGenProjectionTargetMetricsTable',
                           'bestGatesetErrGenProjectionSelfMetricsTable',
                           'logLErrgenProjectionTable',
                           'hamiltonianProjectorTable',
                           'stochasticProjectorTable',
                           'gaugeOptGatesetsVsTargetTable',
                           'gaugeOptCPTPGatesetChoiTable')
        appendix_ls_and_germs_tables = ('byGermTable',)

        tables_to_blank = []
        tables_to_compute = std_tables

        if self._LsAndGermInfoSet:
            tables_to_compute += ls_and_germs_tables
        else:
            tables_to_blank += ls_and_germs_tables

        if showAppendix:
            tables_to_compute += appendix_tables
            if self._LsAndGermInfoSet:
                tables_to_compute += appendix_ls_and_germs_tables
            else:
                tables_to_blank += appendix_ls_and_germs_tables
        else:
            tables_to_blank += appendix_tables
            tables_to_blank += appendix_ls_and_germs_tables


        #Change to report directory so figure generation works correctly
        cwd = _os.getcwd()
        if bOutputFiles and len(report_dir) > 0: _os.chdir(report_dir)

        for key in tables_to_compute:
            qtys[key] = self.tables.get(key, verbosity=printer - 1).render(
                'latex',longtables=self.options.long_tables, scratchDir=D,
                precision=self.options.precision,
                polarprecision=self.options.polar_precision,
                sciprecision=self.options.sci_precision)
            qtys["tt_"+key] = tooltiptex(".tables['%s']" % key)

        _os.chdir(cwd) #change back to original directory

        for key in tables_to_blank:
            qtys[key] = _generation.get_blank_table().render(
                'latex',longtables=self.options.long_tables)
            qtys["tt_"+key] = ""

        #get appendix tables if needed
        #if gaugeOptAppendix:
        #    goaTables = self._specials.get('gaugeOptAppendixTables',verbosity=v)
        #    qtys.update( { key : goaTables[key].render('latex')
        #                   for key in goaTables }  )
        #    #TODO: tables[ref] and then tooltips?
        #
        #elif any((debugAidsAppendix, pixelPlotAppendix, whackamoleAppendix)):
        #    goaTables = self._specials.get('blankGaugeOptAppendixTables',
        #                      verbosity=v)   # fill keys with blank tables
        #    qtys.update( { key : goaTables[key].render('latex')
        #                   for key in goaTables }  )  # for format substitution
        #    #TODO: tables[ref] and then tooltips?


        # 2) generate plots
        printer.log("*** Generating plots *** (%.1fs elapsed)"
                    % (_time.time()-tStart))

        if _matplotlib.is_interactive():
            _matplotlib.pyplot.ioff()
            bWasInteractive = True
        else: bWasInteractive = False

        maxW,maxH = 6.5,8.5 #max width and height of graphic in latex document (in inches)

        def incgr(figFilenm,W=None,H=None): #includegraphics "macro"
            if W is None: W = maxW
            if H is None: H = maxH
            return "\\includegraphics[width=%.2fin,height=%.2fin" % (W,H) + \
                ",keepaspectratio]{%s/%s}" % (D,figFilenm)

        def set_fig_qtys(figkey, figFilenm, v, W=None,H=None):
            fig = self.figures.get(figkey, verbosity=v)
            if bOutputFiles:
                fig.save_to(_os.path.join(report_dir, D, figFilenm))
            qtys[figkey] = incgr(figFilenm,W,H)
            qtys['tt_' + figkey] = tooltiptex(".figures['%s']" % figkey)
            return fig

        #Chi2 or logl plots
        if self._LsAndGermInfoSet:
            Ls = self.parameters['max length list']
            st = 1 if Ls[0] == 0 else 0 #start index: skip LGST column in plots

            if self.parameters['objective'] == "chi2":
                plotFnName,plotFnLatex = "Chi2", "$\chi^2$"
            elif self.parameters['objective'] == "logl":
                plotFnName,plotFnLatex = "LogL", "$\\log(\\mathcal{L})$"
            else:
                raise ValueError("Invalid objective value: %s"
                                 % self.parameters['objective'])

            printer.log(" -- %s plots: " % (plotFnName), end='')
            with printer.progress_logging(1):
                printer.show_progress(0, 3, prefix='', end='')

                w = min(len(self.gatestring_lists['prep fiducials']) * 0.3,maxW)
                h = min(len(self.gatestring_lists['effect fiducials']) * 0.3,maxH)

                fig = set_fig_qtys("colorBoxPlotKeyPlot",
                                   "colorBoxPlotKey.png", printer - 1, w,h)

                printer.show_progress(1, 3, prefix='', end='')

                fig = set_fig_qtys("bestEstimateSummedColorBoxPlot",
                                   "best%sBoxesSummed.png" % plotFnName,
                                   printer - 1,
                                   maxW, maxH-1.0) # -1 for room for caption

                printer.show_progress(2, 3, prefix='', end='')

            figkey = "bestEstimateColorBoxPlotPages"
            figs = self._specials.get(figkey, verbosity=printer - 1)
            incgr_list = []
            for iFig,fig in enumerate(figs):
                figFilenm = "best%sBoxes_pg%d.png" % (plotFnName,iFig)
                if bOutputFiles:
                    fig.save_to(_os.path.join(report_dir, D, figFilenm))
                if iFig == 0:
                    maxX = fig.get_extra_info()['nUsedXs']
                    maxY = fig.get_extra_info()['nUsedYs']
                    incgr_list.append(incgr(figFilenm,maxW,maxH-1.25),)
                else:
                    lx = fig.get_extra_info()['nUsedXs']
                    ly = fig.get_extra_info()['nUsedYs']

                    #scale figure size according to number of rows and columns+1
                    # (+1 for labels ~ another col) relative to initial plot
                    W = float(lx+1)/float(maxX+1) * maxW
                    H = float(ly)  /float(maxY)   * (maxH - 1.25) # -1 for caption
                    incgr_list.append(incgr(figFilenm,W,H))
            qtys[figkey] = "\\end{center}\\end{figure}\\begin{figure}\\begin{center}".join(
                           incgr_list)
            qtys['tt_' + figkey] = tooltiptex(".figures['%s']" % "bestEstimateColorBoxPlot")

        else:
            for figkey in ["colorBoxPlotKeyPlot",
                           "bestEstimateColorBoxPlot"]:
                qtys[figkey] = qtys["tt_"+figkey] = ""
                # "invertedBestEstimateColorBoxPlot"


        pixplots = ""

        #Set template quantity (empty string if appendix disabled)
        qtys['intermediate_pixel_plot_figures'] = pixplots

        printer.log('')

        whackamoleplots = ""

        #Set template quantity (empty string if appendix disabled)
        qtys['whackamole_plot_figures'] = whackamoleplots

        if bWasInteractive:
            _matplotlib.pyplot.ion()


        # 3) populate template latex file => report latex file
        if bOutputFiles:
            printer.log("*** Merging into template file *** (%.1fs elapsed)"
                        % (_time.time()-tStart))
    
            mainTexFilename = _os.path.join(report_dir, report_base + ".tex")
            appendicesTexFilename = _os.path.join(report_dir, report_base + "_appendices.tex")
            pdfFilename = _os.path.join(report_dir, report_base + ".pdf")
    
            mainTemplate = "report_general_main.tex"
            self._merge_template(qtys, mainTemplate, mainTexFilename)
    
    
            # 4) compile report latex file into PDF
            printer.log("Latex file(s) successfully generated.  Attempting to compile with pdflatex...")
            cwd = _os.getcwd()
            if len(report_dir) > 0:
                _os.chdir(report_dir)
    
            try:
                self._compile_latex_report(report_dir, report_base,
                                           self.options.latex_call, printer)
            except _subprocess.CalledProcessError as e:
                printer.error("pdflatex returned code %d " % e.returncode +
                              "Check %s.log to see details." % report_base)
            finally:
                _os.chdir(cwd)
    
            printer.log("Report generation complete! [total time %.0fs]" \
                            % (_time.time()-tStart))
        else:
            printer.log("'Done! (filename is None, so no output files generated)")
            
        return


    #FUTURE?
    #def create_brief_html_page(self, tableClass):
    #    pass
    #
    #def create_full_html_page(self, tableClass):
    #    pass



class ResultOptions(object):
    """ Class encapsulating the display options of a Results instance """
    def __init__(self):
        self.long_tables = False
        self.table_class = "pygstiTbl"
        self.precision = 4
        self.polar_precision = 3
        self.sci_precision = 0
        self.template_path = "."
        self.latex_cmd = "pdflatex"
        # Don't allow LaTeX to try and recover from errors interactively.
        self.latex_opts = ["-interaction=nonstopmode", "-halt-on-error", "-shell-escape"]
        self.latex_call = [self.latex_cmd] + self.latex_opts
        self.errgen_type = "logTiG" #"logG-logT"
        if _os.path.exists("/dev/null"):
            self.latex_postcmd = "-halt-on-error </dev/null >/dev/null"
        else:
            self.latex_postcmd = "" #no /dev/null, so probably not Unix,
                                    #so don't assume halt-on-error works either.

    def describe(self,prefix):
        s = ""
        s += prefix + ".long_tables    -- long latex tables?  %s\n" \
            % str(self.long_tables)
        s += prefix + ".table_class    -- HTML table class = %s\n" \
            % str(self.table_class)
        s += prefix + ".precision      -- precision = %s\n" \
            % str(self.precision)
        s += prefix + ".polar_precision -- precision for polar exponent = %s\n" \
            % str(self.polar_precision)
        s += prefix + ".sci_precision -- precision for scientific notn = %s\n" \
            % str(self.sci_precision)
        s += prefix + ".errgen_type -- type of error generator = %s\n" \
            % str(self.errgen_type)
        s += prefix + ".template_path  -- pyGSTi templates path = '%s'\n" \
            % str(self.template_path)
        s += prefix + ".latex_cmd      -- latex compiling command = '%s'\n" \
            % str(self.latex_cmd)
        s += prefix + ".latex_postcmd  -- latex compiling command postfix = '%s'\n" \
            % str(self.latex_postcmd)
        return s


    def __str__(self):
        s  = "Display options:\n"
        s += self.describe("  ")
        return s

    def copy(self):
        """ Copy this ResultOptions object """
        cpy = ResultOptions()
        cpy.__dict__.update(self.__dict__)
        return cpy


def _to_pdfinfo(list_of_keyval_tuples):

    def sanitize(val):
        if type(val) in (list,tuple):
            sanitized_val = "[" + ", ".join([sanitize(el)
                                             for el in val]) + "]"
        elif type(val) in (dict,_collections.OrderedDict):
            sanitized_val = "Dict[" + \
                ", ".join([ "%s: %s" % (sanitize(k),sanitize(v)) for k,v
                            in val.items()]) + "]"
        elif isinstance(val, _objs.GateSet):
            sanitized_val = "GATESET_DATA"
        else:
            sanitized_val = sanitize_str( str(val) )
        return sanitized_val

    def sanitize_str(s):
        ret = s.replace("^","")
        ret = ret.replace("(","[")
        ret = ret.replace(")","]")
        return ret

    def sanitize_key(s):
        #More stringent string replacement for keys
        ret = s.replace(" ","_")
        ret = ret.replace("^","")
        ret = ret.replace(",","_")
        ret = ret.replace("(","[")
        ret = ret.replace(")","]")
        return ret


    sanitized_list = []
    for key,val in list_of_keyval_tuples:
        sanitized_key = sanitize_key(key)
        sanitized_val = sanitize(val)
        sanitized_list.append( (sanitized_key, sanitized_val) )

    return ",\n".join( ["%s={%s}" % (key,val) for key,val in sanitized_list] )
